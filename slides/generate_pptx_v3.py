#!/usr/bin/env python3
"""
Generate Workshop_v3.pptx — ESP32-CAM AI Vision Workshop (RMUTT)
v3: dark-navy palette (v1) + Sarabun font + modern layout + real images
Usage : python generate_pptx_v3.py   (run from inside slides/)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

IMG = Path(__file__).parent   # image files live alongside this script

# ── Palette (v1 dark navy) ────────────────────────────────────────────────────
BG     = RGBColor(0x0D, 0x1B, 0x2A)
BG2    = RGBColor(0x0A, 0x13, 0x1E)
GHOST  = RGBColor(0x17, 0x2A, 0x3E)
CYAN   = RGBColor(0x00, 0xC8, 0xFF)
ORANGE = RGBColor(0xFF, 0x70, 0x43)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xD0, 0xE4, 0xF5)
GRAY   = RGBColor(0x88, 0xA0, 0xB8)
GREEN  = RGBColor(0x4C, 0xD9, 0x64)
SEP    = RGBColor(0x1E, 0x32, 0x48)   # subtle row separator

# ── Fonts ─────────────────────────────────────────────────────────────────────
FH = "Sarabun"   # headings
FB = "Sarabun"   # body
FC = "Consolas"  # code

# ── Slide dimensions 16:9 ─────────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
#  PRIMITIVES
# ═══════════════════════════════════════════════════════════════════════════════

def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def _tb(slide, text, left, top, width, height,
        size=18, color=LIGHT, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True, font=FB):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return tb


def _box(slide, left, top, width, height, fill_color, border=False, border_color=None):
    s = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border:
        s.line.color.rgb = border_color or CYAN
    else:
        s.line.fill.background()
    return s


def _oval(slide, left, top, width, height, fill_color):
    s = slide.shapes.add_shape(
        9, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def _multiline(slide, lines, left, top, width, height,
               size=18, color=LIGHT, bold=False,
               align=PP_ALIGN.LEFT, gap=8, font=FB):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, col, sz, bd = item, color, size, bold
        elif len(item) == 2:
            text, col = item; sz, bd = size, bold
        elif len(item) == 3:
            text, col, sz = item; bd = bold
        else:
            text, col, sz, bd = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.bold = bd
        r.font.name = font
    return tb


def _photo(slide, filename, left, top, width, height, caption=""):
    """Insert image; fall back to styled placeholder if file missing."""
    path = IMG / filename
    if path.exists():
        slide.shapes.add_picture(
            str(path), Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        _box(slide, left, top, width, height, BG2, border=True)
        _tb(slide, "📷", left, top + (height - 0.6) / 2, width, 0.55,
            size=38, color=GHOST, align=PP_ALIGN.CENTER)
        if caption:
            _tb(slide, caption, left, top + height - 0.42, width, 0.38,
                size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


def _code_box(slide, code, left, top, width, height):
    _box(slide, left, top, width, height, BG2, border=True)
    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12),
        Inches(width - 0.3), Inches(height - 0.25))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.size = Pt(14)
        r.font.color.rgb = GREEN
        r.font.name = FC


# ── Compound elements ─────────────────────────────────────────────────────────

def _accent_bar(slide, y=7.25, color=CYAN):
    _box(slide, 0, y, 13.33, 0.08, color)


def _top_bar(slide, color=CYAN):
    _box(slide, 0, 0, 13.33, 0.07, color)


def _section_ghost(slide, num):
    _tb(slide, str(num), 0.1, 0.5, 4.5, 6.5,
        size=220, color=GHOST, bold=True, align=PP_ALIGN.LEFT, font=FH)


def _lab_badge(slide, num, color=ORANGE):
    _box(slide, 11.0, 0.18, 2.1, 0.62, color)
    _tb(slide, f"LAB  {num:02d}", 11.0, 0.18, 2.1, 0.62,
        size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FH)


def _slide_title(slide, text, sub=None):
    _tb(slide, text, 0.5, 0.3, 12.3, 0.95,
        size=36, color=WHITE, bold=True, font=FH)
    _box(slide, 0.5, 1.18, 2.5, 0.05, CYAN)
    if sub:
        _tb(slide, sub, 0.5, 1.28, 12.3, 0.48,
            size=15, color=GRAY, italic=True, font=FB)


def _step_badge(slide, step_num, label="", color=ORANGE):
    """Oval number badge + label (modern, no sharp corners)."""
    _oval(slide, 0.5, 1.58, 0.74, 0.74, color)
    _tb(slide, str(step_num), 0.5, 1.58, 0.74, 0.74,
        size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FH)
    if label:
        _tb(slide, label, 1.38, 1.62, 11.3, 0.6,
            size=21, color=WHITE, bold=True, font=FH)


def _checklist(slide, items, left=0.8, top=2.0, size=19):
    lines = [(f"  {item}", WHITE, size, False) for item in items]
    # Draw oval bullet for each item
    y = top
    for item in items:
        _oval(slide, left, y + 0.08, 0.32, 0.32, CYAN)
        _tb(slide, "✓", left, y + 0.04, 0.32, 0.38,
            size=13, color=BG, bold=True, align=PP_ALIGN.CENTER, font=FH)
        _tb(slide, item, left + 0.45, y, 11.4, 0.48,
            size=size, color=LIGHT, font=FB)
        y += 0.82


def _list_row(slide, title, desc, left, top, width=12.3, accent=CYAN, title_color=WHITE):
    """Modern: thin vertical accent line + title + description. No box."""
    _box(slide, left, top + 0.05, 0.05, 0.85, accent)
    _tb(slide, title, left + 0.22, top + 0.04, width - 0.22, 0.42,
        size=18, color=title_color, bold=True, font=FH)
    if desc:
        _tb(slide, desc, left + 0.22, top + 0.46, width - 0.22, 0.38,
            size=14, color=LIGHT, font=FB)


def _info_row(slide, label, value, top, label_color=CYAN):
    """Two-column info row used in Gemini/Claude spec tables."""
    _tb(slide, label, 0.8, top, 2.5, 0.48,
        size=16, color=label_color, bold=True, font=FH)
    _tb(slide, value, 3.4, top, 9.2, 0.48,
        size=16, color=LIGHT, font=FB)
    _box(slide, 0.8, top + 0.5, 11.8, 0.01, SEP)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 0 — OPENING
# ═══════════════════════════════════════════════════════════════════════════════

def s00_cover(prs):
    slide = _new_slide(prs)

    # Hero image — right panel
    _photo(slide, "Cover Slide Hero (Slide 01) - v1.jpeg", 6.7, 0, 6.63, 7.5)

    # Left panel accents
    _box(slide, 0, 0,    6.7, 0.12, CYAN)
    _box(slide, 0, 7.38, 6.7, 0.12, CYAN)
    _box(slide, 0, 0.12, 0.08, 7.26, ORANGE)

    _tb(slide, "ESP32-CAM", 0.5, 1.1, 6.0, 1.0,
        size=54, color=CYAN, bold=True, font=FH)
    _tb(slide, "AI Vision Workshop", 0.5, 2.0, 6.0, 0.85,
        size=36, color=WHITE, bold=True, font=FH)

    _box(slide, 0.5, 3.0, 4.0, 0.05, ORANGE)

    _tb(slide, "การพัฒนากล้องอัจฉริยะด้วย ESP32-CAM", 0.5, 3.18, 6.0, 0.5,
        size=16, color=LIGHT, font=FB)
    _tb(slide, "และปัญญาประดิษฐ์สำหรับการวิเคราะห์ภาพ", 0.5, 3.65, 6.0, 0.5,
        size=16, color=LIGHT, font=FB)

    _tb(slide, "มหาวิทยาลัยเทคโนโลยีราชมงคลธัญบุรี (RMUTT)",
        0.5, 4.8, 6.0, 0.45, size=13, color=GRAY, font=FB)
    _tb(slide, "[วันที่ / Date]  •  #ESP32CAM_AI_RMUTT",
        0.5, 5.25, 6.0, 0.4, size=12, color=GRAY, font=FB)


def s00_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Workshop Overview", "ในวันนี้เราจะสร้างอะไร?")

    _tb(slide, "เราจะสร้างระบบกล้อง AI ที่สามารถ...",
        0.5, 1.82, 12.5, 0.48, size=20, color=CYAN, bold=True, font=FH)
    _tb(slide, "วิเคราะห์ภาพ  •  นับวัตถุ  •  จำแนกประเภท  •  ตอบคำถาม",
        0.5, 2.28, 12.5, 0.42, size=17, color=LIGHT, font=FB)

    _box(slide, 0.5, 2.85, 12.3, 3.95, BG2, border=True)
    flow = [
        ("📷  ESP32-CAM",                       CYAN,   19, True),
        ("        │  Wi-Fi  →  /capture",        GRAY,   15, False),
        ("        ▼",                             GRAY,   14, False),
        ("⚙️   FastAPI Backend  (Python)",        LIGHT,  17, False),
        ("        │  image_bytes + prompt",       GRAY,   15, False),
        ("        ▼",                             GRAY,   14, False),
        ("🤖  Vision AI  (Gemini / Claude)",      ORANGE, 17, True),
        ("        │  JSON response",              GRAY,   15, False),
        ("        ▼",                             GRAY,   14, False),
        ("📊  Web UI  —  ผลการวิเคราะห์ภาพ",     GREEN,  17, True),
    ]
    _multiline(slide, flow, 3.8, 2.95, 7.0, 3.8, gap=2)


def s00_outcomes(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Learning Outcomes", "เมื่อสิ้นสุด Workshop นักเรียนจะสามารถ:")

    outcomes = [
        ("1.  ติดตั้งและใช้งาน ESP32-CAM",            "Configure & Flash Firmware"),
        ("2.  อธิบาย REST API",                         "Explain REST API & HTTP"),
        ("3.  รัน Python Backend ด้วย FastAPI",         "Run FastAPI Server"),
        ("4.  ใช้ Vision AI API วิเคราะห์ภาพ",          "Use Gemini / Claude Vision"),
        ("5.  เขียน Prompt ที่เจาะจง",                  "Design Effective Prompts"),
        ("6.  สร้าง AI Vision Mini Project",             "Build a Real AI Project"),
    ]
    top = 1.72
    for th, en in outcomes:
        _box(slide, 0.5, top, 0.05, 0.72, CYAN)
        _tb(slide, th, 0.75, top + 0.04, 12.0, 0.4, size=19, color=WHITE, bold=True, font=FH)
        _tb(slide, en, 0.75, top + 0.42, 12.0, 0.28, size=12, color=GRAY, italic=True, font=FB)
        top += 0.82


def s00_agenda(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Workshop Agenda", "ตารางเวลา 1 วัน")

    rows = [
        ("08:30–09:00", "ลงทะเบียน / รับอุปกรณ์",            False),
        ("09:00–09:30", "Lecture: AI + IoT + Computer Vision",  True),
        ("09:30–10:30", "Lab 01 — ESP32-CAM Camera Server",     True),
        ("10:30–10:45", "พัก ☕",                               False),
        ("10:45–11:45", "Lab 02 — REST API & Backend",          True),
        ("11:45–12:45", "พักกลางวัน 🍱",                       False),
        ("12:45–13:00", "Lecture: Prompt Engineering",           True),
        ("13:00–14:15", "Lab 03 — AI Vision Analysis",          True),
        ("14:15–14:30", "พัก",                                  False),
        ("14:30–15:45", "Lab 04 — Mini Project",                True),
        ("15:45–16:00", "นำเสนอผลงาน / มอบใบประกาศ",          False),
    ]
    top = 1.58
    for time, activity, highlight in rows:
        col = CYAN if highlight else LIGHT
        bd  = highlight
        _tb(slide, time,     0.4, top, 2.3, 0.38, size=13, color=GRAY, font=FB)
        if highlight:
            _box(slide, 2.75, top + 0.08, 0.04, 0.26, CYAN)
        _tb(slide, activity, 2.9, top, 10.0, 0.38, size=14, color=col, bold=bd, font=FB if not highlight else FH)
        top += 0.38


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION HEADER — shared (with background image on right)
# ═══════════════════════════════════════════════════════════════════════════════

def _section_header(prs, num, title_en, title_th, accent=CYAN):
    slide = _new_slide(prs)

    # Right panel: background image
    _photo(slide, "Section Header Background (Section 1-4).jpeg",
           6.83, 0, 6.5, 7.5)
    # Gradient edge (overlap strip so text side stays cleanly dark)
    _box(slide, 6.5, 0, 0.6, 7.5, BG)

    _section_ghost(slide, num)
    _box(slide, 0, 0, 0.12, 7.5, accent)
    _box(slide, 0.12, 6.82, 6.4, 0.06, accent)

    _tb(slide, title_en, 1.3, 2.0, 5.4, 1.3,
        size=48, color=WHITE, bold=True, font=FH)
    _tb(slide, title_th, 1.3, 3.35, 5.4, 0.72,
        size=20, color=accent, font=FH)
    _tb(slide, f"Section {num}", 1.3, 4.12, 5.4, 0.4,
        size=14, color=GRAY, font=FB)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 1 — IoT & ESP32-CAM
# ═══════════════════════════════════════════════════════════════════════════════

def s01_header(prs):
    _section_header(prs, 1, "IoT & ESP32-CAM", "อินเทอร์เน็ตของสรรพสิ่ง และกล้อง AI")


def s01_iot(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "What is IoT?", "Internet of Things — อินเทอร์เน็ตของสรรพสิ่ง")

    _tb(slide, "IoT = อุปกรณ์ที่เชื่อมต่ออินเทอร์เน็ตเพื่อส่ง/รับข้อมูล",
        0.5, 1.82, 12.5, 0.52, size=21, color=CYAN, bold=True, font=FH)

    examples = [
        ("🏠  Smart Home",        "หลอดไฟ, กล้องวงจรปิด, ล็อคประตู ที่ควบคุมผ่านมือถือ"),
        ("🏭  Smart Factory",     "เซนเซอร์ตรวจจับอุณหภูมิ, แรงสั่นสะเทือน, แจ้งเตือนอัตโนมัติ"),
        ("🌱  Smart Agriculture", "กล้องถ่ายภาพต้นไม้ → AI วิเคราะห์โรค → แจ้งเตือนเกษตรกร"),
    ]
    top = 2.52
    for title, desc in examples:
        _list_row(slide, title, desc, 0.5, top, accent=CYAN)
        top += 1.38


def s01_esp32cam(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "ESP32-CAM Introduction", "โมดูลกล้อง AI Thinker — หัวใจของ Workshop")

    specs = [
        ("CPU",    "ESP32 Dual-Core 240 MHz"),
        ("RAM",    "520 KB SRAM  +  4 MB PSRAM"),
        ("Flash",  "4 MB"),
        ("Camera", "OV2640  2MP  (สูงสุด 1600×1200)"),
        ("Wi-Fi",  "802.11 b/g/n  2.4 GHz"),
        ("Price",  "~150–250 บาท"),
    ]
    top = 1.72
    for label, val in specs:
        _tb(slide, label, 0.6, top, 2.2, 0.44, size=15, color=CYAN, bold=True, font=FH)
        _tb(slide, val,   2.9,  top, 5.0, 0.44, size=15, color=LIGHT, font=FB)
        _box(slide, 0.6, top + 0.46, 7.2, 0.01, SEP)
        top += 0.62

    # Hardware board photo (right panel)
    _photo(slide, "ESP32-CAM Board Diagram (Slide 07).png",
           8.1, 1.55, 4.9, 5.55, "AI Thinker ESP32-CAM")


def s01_hardware(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "ESP32-CAM Hardware Tour", "รู้จักส่วนประกอบสำคัญ")

    diagram = """\
  ┌─────────────────────────────┐
  │        OV2640 Camera        │  ← ถอด-เสียบได้
  │  ┌───────────────────────┐  │
  │  │    ESP32-WROVER SoC   │  │  ← CPU + Wi-Fi + RAM
  │  └───────────────────────┘  │
  │  GPIO0 ── ต่อ GND ขณะ flash │  ← ⚠️ สำคัญมาก
  │  5V / GND ── จ่ายไฟ         │
  │  GPIO1/3  ── UART (Serial)  │
  └─────────────────────────────┘
         │  USB (ผ่าน MB Board)
         ▼
     ESP32-CAM-MB  (CH340)"""
    _code_box(slide, diagram, 0.5, 1.72, 12.3, 5.4)


def s01_wifi(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Wi-Fi & Networking Basics", "เครือข่ายในห้อง Workshop")

    points = [
        ("SSID",        "ชื่อ Wi-Fi ที่มองเห็น  → ใส่ใน secrets.h"),
        ("Password",    "รหัส Wi-Fi             → ใส่ใน secrets.h"),
        ("IP Address",  "เช่น 192.168.1.50 — ที่อยู่ของ ESP32-CAM ในเครือข่าย"),
        ("Port",        "Web UI: port 80  •  Stream: port 81"),
        ("⚠️  2.4 GHz", "ESP32 รองรับแค่ 2.4 GHz  —  ไม่รองรับ 5 GHz"),
    ]
    top = 1.78
    for label, desc in points:
        _list_row(slide, label, desc, 0.5, top, accent=CYAN)
        top += 1.08


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 2 — REST API & BACKEND
# ═══════════════════════════════════════════════════════════════════════════════

def s02_header(prs):
    _section_header(prs, 2, "REST API & Backend", "การสื่อสารระหว่างอุปกรณ์กับ AI")


def s02_api(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "What is REST API?", "API คืออะไร? — เปรียบเหมือนร้านอาหาร")

    cols = [
        ("🍽️  ร้านอาหาร", ["ลูกค้า = Client", "พนักงาน = API", "ครัว = Server", "เมนู = Endpoints", "ออเดอร์ = Request", "อาหาร = Response"]),
        ("💻  REST API",   ["Browser = Client", "API = Interface", "FastAPI = Server", "GET /prompts, POST /analyze", "HTTP Request + JSON", "JSON Response"]),
    ]
    for i, (heading, items) in enumerate(cols):
        left = 0.5 + i * 6.55
        _box(slide, left, 1.75, 6.3, 5.35, BG2, border=True)
        _tb(slide, heading, left + 0.22, 1.85, 5.9, 0.5,
            size=19, color=CYAN, bold=True, font=FH)
        _box(slide, left, 2.3, 6.3, 0.04, CYAN)
        for j, item in enumerate(items):
            _tb(slide, item, left + 0.22, 2.48 + j * 0.68, 5.9, 0.55,
                size=16, color=LIGHT, font=FB)


def s02_http(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "HTTP GET vs POST", "วิธีส่งข้อมูลผ่าน HTTP")

    for i, (method, color, use, example) in enumerate([
        ("GET",  CYAN,   "ขอดึงข้อมูล (Read)",       "GET /prompts\nGET /health"),
        ("POST", ORANGE, "ส่งข้อมูลไป (Write/Send)", "POST /analyze\n{ cam_url, prompt }"),
    ]):
        left = 0.5 + i * 6.55
        _box(slide, left, 1.75, 6.3, 5.4, BG2, border=True)
        _box(slide, left, 1.75, 6.3, 0.65, color)
        _tb(slide, method, left + 0.22, 1.78, 5.9, 0.58,
            size=28, color=WHITE, bold=True, font=FH)
        _tb(slide, use, left + 0.22, 2.58, 5.9, 0.48,
            size=17, color=LIGHT, bold=True, font=FB)
        _code_box(slide, example, left + 0.22, 3.18, 5.88, 1.52)
        _tb(slide, "ไม่มี body" if i == 0 else "ส่ง body เป็น JSON",
            left + 0.22, 4.85, 5.9, 0.38,
            size=13, color=GRAY, italic=True, font=FB)


def s02_json(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "JSON Basics", "รูปแบบข้อมูลที่ใช้สื่อสารระหว่าง Backend กับ AI")

    _tb(slide, "JSON = JavaScript Object Notation — ข้อมูล key: value",
        0.5, 1.78, 12.5, 0.48, size=16, color=CYAN, font=FB)

    sample = """\
{
  "provider": "gemini",
  "filename":  "photo.jpg",
  "image_size_bytes": 12480,
  "result": "ภาพแสดงโต๊ะทำงานที่มีโน้ตบุ๊ค 1 เครื่อง
             กาแฟ 1 แก้ว และหนังสือ 3 เล่ม"
}"""
    _code_box(slide, sample, 0.5, 2.35, 7.5, 4.72)

    _multiline(slide, [
        ("Key Points:", CYAN, 16, True),
        ("• ใช้  { }  สำหรับ object", LIGHT, 15, False),
        ("• ใช้  [ ]  สำหรับ list/array", LIGHT, 15, False),
        ("• Key ใส่ใน  \" \"  เสมอ", LIGHT, 15, False),
        ("• Value: string, number,", LIGHT, 15, False),
        ("  boolean, list, object", LIGHT, 15, False),
        ("• ทุก API ตอบกลับเป็น JSON", LIGHT, 15, False),
    ], 8.3, 2.35, 4.8, 4.72, gap=10, font=FB)


def s02_fastapi(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "FastAPI Backend Overview", "สถาปัตยกรรมของ Backend")

    arch = """\
Browser / Web UI  (localhost:8000)
        │
        │  POST /analyze   { cam_url, prompt }
        ▼
  ┌─────────────────────────────────────┐
  │           FastAPI (main.py)         │
  │   GET /       → index.html          │
  │   GET /prompts → preset list        │
  │   POST /analyze → fetch + AI        │
  │   POST /analyze-upload → AI         │
  │   GET /health  → status             │
  └──────────────┬──────────────────────┘
                 │                  │
        ┌────────┴────────┐  ┌──────┴──────┐
        │  ESP32-CAM       │  │  Vision AI  │
        │  /capture (JPEG) │  │  mock/gemini│
        └──────────────────┘  │  /claude    │
                              └─────────────┘"""
    _code_box(slide, arch, 0.5, 1.72, 12.3, 5.45)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 3 — AI VISION
# ═══════════════════════════════════════════════════════════════════════════════

def s03_header(prs):
    _section_header(prs, 3, "AI & Vision", "ปัญญาประดิษฐ์และการวิเคราะห์ภาพ", ORANGE)


def s03_ai(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "What is AI?", "ปัญญาประดิษฐ์คืออะไร?")

    stages = [
        ("Rule-Based",    "กฎตายตัว\nIF วัตถุกลม AND สีแดง → THEN แอปเปิ้ล", GRAY),
        ("Machine\nLearning", "เรียนรู้จากตัวอย่าง\n1,000 รูป → เรียนรู้เอง → ทำนาย", CYAN),
        ("Generative AI", "สร้างคำตอบใหม่\nรับรูป + prompt → ตอบภาษาธรรมชาติ", ORANGE),
    ]
    for i, (title, desc, color) in enumerate(stages):
        left = 0.5 + i * 4.27
        _box(slide, left, 1.75, 4.0, 5.35, BG2, border=True)
        _box(slide, left, 1.75, 4.0, 0.55, color)
        _tb(slide, title, left + 0.15, 1.78, 3.7, 0.5,
            size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FH)
        _tb(slide, desc, left + 0.15, 2.55, 3.7, 4.3,
            size=15, color=LIGHT, wrap=True, font=FB)
        if i < 2:
            _tb(slide, "→", left + 4.08, 3.9, 0.4, 0.48,
                size=22, color=CYAN, bold=True, font=FH)


def s03_genai(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Generative AI", "AI ที่สร้างคำตอบใหม่จากข้อมูลที่ได้รับ")

    items = [
        ("📝  Text Generation",  "เขียนเรียงความ, สรุปข้อมูล, ตอบคำถาม"),
        ("🖼️  Image Generation", "สร้างภาพจากคำอธิบาย (Midjourney, DALL-E)"),
        ("👁️  Image Analysis",   "อธิบาย, จำแนก, นับ สิ่งที่อยู่ในภาพ  ← Workshop นี้!"),
        ("💻  Code Generation",  "เขียนโปรแกรมจาก requirements"),
    ]
    top = 1.85
    for title, desc in items:
        highlight = "Workshop" in desc
        accent = GREEN if highlight else CYAN
        _list_row(slide, title, desc, 0.5, top, accent=accent, title_color=GREEN if highlight else WHITE)
        top += 1.28


def s03_vision(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Vision AI", "AI ที่ 'มองเห็น' และเข้าใจภาพ")

    _tb(slide, "Vision AI รับภาพ + prompt แล้วตอบกลับเป็นข้อความ",
        0.5, 1.72, 12.5, 0.48, size=19, color=CYAN, font=FH)

    capabilities = [
        ("🔍 Object Detection",   "ระบุวัตถุในภาพ: แมว, โต๊ะ, ขวด"),
        ("🔢 Counting",            "นับจำนวน: มีหนังสือ 3 เล่ม"),
        ("🎨 Color / Scene",      "วิเคราะห์สีและบรรยากาศ"),
        ("📊 Classification",      "จำแนกประเภท: ขยะรีไซเคิล"),
        ("💚 Health Assessment",  "ประเมินสุขภาพ: ต้นไม้มีใบเหลือง"),
        ("📝 Text Reading (OCR)", "อ่านตัวหนังสือในภาพ"),
    ]
    top = 2.35
    for i, (cap, example) in enumerate(capabilities):
        col = 0 if i % 2 == 0 else 6.6
        row = top + (i // 2) * 1.38
        _box(slide, 0.5 + col, row, 6.3, 1.15, BG2, border=True)
        _tb(slide, cap, 0.75 + col, row + 0.08, 5.8, 0.4,
            size=16, color=ORANGE, bold=True, font=FH)
        _tb(slide, example, 0.75 + col, row + 0.52, 5.8, 0.45,
            size=14, color=LIGHT, font=FB)


def s03_gemini(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Gemini Vision (Google)", "Vision AI โดย Google — แนะนำสำหรับ Workshop")

    _box(slide, 0.5, 1.72, 12.3, 5.45, BG2, border=True)
    info = [
        ("Model",     "gemini-2.5-flash"),
        ("Provider",  "Google AI Studio"),
        ("API Key",   "ฟรีสำหรับนักเรียน — aistudio.google.com"),
        ("Free Tier", "1,500 requests/วัน  (เพียงพอสำหรับ Workshop)"),
        ("SDK",       "google-genai  (python)"),
        ("ENV",       "VISION_PROVIDER=gemini"),
    ]
    top = 1.92
    for label, val in info:
        _info_row(slide, label, val, top, CYAN)
        top += 0.75

    _tb(slide, "✅  แนะนำให้เริ่มด้วย Gemini เพราะ free tier สูงที่สุด",
        0.8, 6.6, 12.0, 0.42, size=14, color=GREEN, bold=True, font=FB)


def s03_claude(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Claude Vision (Anthropic)", "Vision AI โดย Anthropic")

    _box(slide, 0.5, 1.72, 12.3, 5.45, BG2, border=True)
    info = [
        ("Model",    "claude-haiku-4-5"),
        ("Provider", "Anthropic"),
        ("API Key",  "console.anthropic.com  (ต้องใส่บัตรเครดิต)"),
        ("Pricing",  "$1.00 / M input tokens  •  $5.00 / M output tokens"),
        ("SDK",      "anthropic  (python)"),
        ("ENV",      "VISION_PROVIDER=claude"),
    ]
    top = 1.92
    for label, val in info:
        _info_row(slide, label, val, top, ORANGE)
        top += 0.75

    _tb(slide, "💡  เปลี่ยนแค่ VISION_PROVIDER=claude ใน .env",
        0.8, 6.6, 12.0, 0.42, size=14, color=CYAN, font=FB)


def s03_prompt(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Prompt Engineering", "วิธีเขียนคำสั่งให้ AI ตอบได้ตรงใจ")

    _box(slide, 0.5, 1.72, 12.3, 0.72, BG2, border=True)
    _tb(slide, "[Role]  +  [Task]  +  [Format]  +  [Language]",
        0.7, 1.77, 12.0, 0.6,
        size=22, color=CYAN, bold=True, align=PP_ALIGN.CENTER, font=FH)

    parts = [
        ("Role",     "คุณคือผู้เชี่ยวชาญด้านสิ่งแวดล้อม",  "บอก AI ว่าเป็นใคร"),
        ("Task",     "วิเคราะห์ขยะในภาพและจำแนกประเภท",    "สิ่งที่ต้องทำ"),
        ("Format",   "ตอบเป็น bullet point",                  "รูปแบบคำตอบ"),
        ("Language", "ภาษาไทย",                              "ภาษาที่ใช้"),
    ]
    top = 2.65
    for label, example, desc in parts:
        _box(slide, 0.5, top, 12.3, 1.02, BG2, border=True)
        _tb(slide, label,            0.75, top + 0.06, 1.8, 0.4, size=16, color=ORANGE, bold=True, font=FH)
        _tb(slide, f'"{example}"',   2.65, top + 0.06, 5.5, 0.4, size=15, color=GREEN, font=FB)
        _tb(slide, desc,             8.3,  top + 0.06, 4.3, 0.4, size=13, color=GRAY, italic=True, font=FB)
        top += 1.18


def s03_workflow(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "AI Vision Workflow", "ข้อมูลไหลอย่างไร จากกล้องถึง AI")

    workflow = """\
  👤 นักเรียน เปิด Web UI (localhost:8000)
        │  กรอก cam_url + prompt → กด Capture & Analyze
        ▼
  ⚙️  FastAPI Backend  (main.py)
        │
        ├─── GET http://<esp32-ip>/capture  ──► 📷 ESP32-CAM
        │         (JPEG image, ~10–50 KB)         └── ถ่ายภาพ → ส่งกลับ
        │
        ├─── image_bytes + prompt
        │         ▼
        │    Vision Provider  (VISION_PROVIDER env var)
        │         ├── mock   → canned response
        │         ├── gemini → Google AI API
        │         └── claude → Anthropic API
        │
        └── JSON Response  ──► 🌐 Web UI แสดงผล"""
    _code_box(slide, workflow, 0.5, 1.72, 12.3, 5.45)


# ═══════════════════════════════════════════════════════════════════════════════
#  LAB HEADER helper
# ═══════════════════════════════════════════════════════════════════════════════

def _lab_header(prs, num, title_en, title_th, time_str, accent=ORANGE):
    slide = _new_slide(prs)

    # Right panel background image
    _photo(slide, "Section Header Background (Section 1-4).jpeg",
           6.83, 0, 6.5, 7.5)
    _box(slide, 6.5, 0, 0.6, 7.5, BG)  # dark fade strip

    _section_ghost(slide, num)
    _box(slide, 0, 0, 0.12, 7.5, accent)
    _box(slide, 0.12, 6.82, 6.4, 0.06, accent)

    _tb(slide, f"Lab {num:02d}", 1.3, 1.5, 5.4, 1.1,
        size=58, color=accent, bold=True, font=FH)
    _tb(slide, title_en, 1.3, 2.6, 5.4, 0.85,
        size=28, color=WHITE, bold=True, font=FH)
    _tb(slide, title_th, 1.3, 3.48, 5.4, 0.6,
        size=17, color=LIGHT, font=FB)
    _tb(slide, f"เวลา: {time_str}", 1.3, 4.2, 4.0, 0.4,
        size=14, color=GRAY, font=FB)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 4 — LAB 01
# ═══════════════════════════════════════════════════════════════════════════════

def s04_header(prs):
    _lab_header(prs, 1, "ESP32-CAM Camera Server",
                "ติดตั้งกล้อง • Flash Firmware • ดู Live Stream", "~60 นาที", ORANGE)


def s04_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Lab 01 — Overview", "สิ่งที่เราจะทำใน Lab นี้")

    steps = [
        "ต่อสาย ESP32-CAM กับ ESP32-CAM-MB board",
        "ติดตั้ง Arduino IDE และ ESP32 Board Package",
        "แก้ไข secrets.h ใส่ Wi-Fi SSID และ Password",
        "Flash firmware เข้า ESP32-CAM",
        "ทดสอบ Live Stream และ /capture endpoint",
    ]
    _checklist(slide, steps, top=1.85)
    _tb(slide, "🎯  เป้าหมาย: กล้อง online และ /capture ส่ง JPEG ได้",
        0.5, 6.72, 12.3, 0.42, size=14, color=GREEN, bold=True, font=FB)


def s04_step1(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 1 — Wire ESP32-CAM")
    _step_badge(slide, 1, "ต่อสาย ESP32-CAM + MB Board")

    # Left: key warnings
    _multiline(slide, [
        ("⚠️  GPIO0 → GND", ORANGE, 19, True),
        ("ต้องทำก่อน Upload ทุกครั้ง", LIGHT, 16, False),
        ("", GRAY, 8, False),
        ("หลัง Upload เสร็จ:", LIGHT, 16, True),
        ("→ ถอด jumper GPIO0 ออก", CYAN, 16, False),
        ("→ กดปุ่ม RST บน MB Board", CYAN, 16, False),
        ("", GRAY, 8, False),
        ("🔌  ESP32-CAM → MB Board → USB → PC", GRAY, 14, False),
    ], 0.5, 2.42, 6.5, 4.6, gap=8, font=FB)

    # Right: wiring photo
    _photo(slide, "Lab 01 Wiring Setup (Slide 25).jpeg",
           7.2, 2.35, 5.8, 4.65, "Wiring Setup")


def s04_step2(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 2 — Arduino IDE Setup")
    _step_badge(slide, 2, "ตั้งค่า Arduino IDE")

    steps_txt = """\
1. เปิด Arduino IDE 2.x
2. File → Preferences
   → Additional boards manager URLs:
   https://dl.espressif.com/dl/package_esp32_index.json

3. Tools → Board → Boards Manager
   → ค้นหา "esp32" → Install (Espressif Systems)

4. Tools → Board → ESP32 Arduino
   → เลือก "AI Thinker ESP32-CAM"

5. Tools → Port → เลือก COMx (ของ MB Board)"""
    _code_box(slide, steps_txt, 0.5, 2.42, 12.3, 4.62)


def s04_step3(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 3 — Edit secrets.h")
    _step_badge(slide, 3, "ใส่ Wi-Fi credentials")

    _tb(slide, "คัดลอก  secrets.example.h  →  secrets.h  แล้วแก้ไข:",
        0.5, 2.42, 12.5, 0.42, size=15, color=LIGHT, font=FB)

    code = """\
// secrets.h
#define WIFI_SSID  "ชื่อ Wi-Fi ของคุณ"
#define WIFI_PASS  "รหัส Wi-Fi ของคุณ"
"""
    _code_box(slide, code, 0.5, 3.0, 12.3, 1.58)

    _multiline(slide, [
        ("⚠️  secrets.h อยู่ใน .gitignore — ไม่ถูก commit ขึ้น GitHub", ORANGE, 15, True),
        ("✅  ใช้ 2.4 GHz เท่านั้น (ESP32 ไม่รองรับ 5 GHz)", GREEN, 15, False),
        ("✅  ตรวจสอบให้ครั้งสุดท้ายก่อน Upload", GREEN, 15, False),
    ], 0.5, 4.75, 12.5, 2.2, gap=10, font=FB)


def s04_step4(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 4 — Flash Firmware")
    _step_badge(slide, 4, "Upload Firmware เข้า ESP32-CAM")

    steps_txt = """\
1. ตรวจสอบ: GPIO0 ต่อลง GND แล้ว
2. กด Upload ใน Arduino IDE (→)
3. รอ "Connecting......____"
   → ถ้าค้าง: ลองกด RST สั้นๆ
4. เห็น "Writing at 0x..." = กำลัง flash ✅
5. เห็น "Hard resetting..." = flash เสร็จ ✅
6. ถอด jumper GPIO0 ออก
7. กด RST อีกครั้ง → ESP32-CAM บูตปกติ"""
    _code_box(slide, steps_txt, 0.5, 2.42, 7.5, 4.62)

    _multiline(slide, [
        ("ปัญหาที่พบบ่อย:", ORANGE, 16, True),
        ("❌ ค้างที่ Connecting", LIGHT, 15, False),
        ("→ GPIO0 ไม่ได้ต่อ GND", GRAY, 14, False),
        ("❌ Upload failed", LIGHT, 15, False),
        ("→ เลือก Port ผิด", GRAY, 14, False),
        ("❌ ไม่เห็น Port", LIGHT, 15, False),
        ("→ ไม่มีไดรเวอร์ CH340", GRAY, 14, False),
    ], 8.3, 2.42, 4.8, 4.62, gap=7, font=FB)


def s04_step5(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 5 — Find IP & Test")
    _step_badge(slide, 5, "หา IP Address และทดสอบ")

    code = """\
# Serial Monitor → Baud Rate: 115200
Connecting to WiFi......
WiFi connected!
IP Address: 192.168.x.x   ← จด IP นี้ไว้!
Camera Ready! Stream: http://192.168.x.x:81/stream"""
    _code_box(slide, code, 0.5, 2.42, 7.5, 2.8)

    _multiline(slide, [
        ("ทดสอบใน Browser:", CYAN, 16, True),
        ("http://<IP>          → Web UI", LIGHT, 15, False),
        ("http://<IP>/capture  → JPEG", LIGHT, 15, False),
        ("http://<IP>:81/stream → Live", LIGHT, 15, False),
    ], 8.3, 2.42, 4.8, 2.8, gap=9, font=FB)

    _tb(slide, "📝  จด IP Address ไว้ใช้ใน Lab 02-03-04",
        0.5, 5.42, 12.3, 0.42, size=14, color=ORANGE, bold=True, font=FB)


def s04_result(prs):
    slide = _new_slide(prs)
    _top_bar(slide, GREEN); _accent_bar(slide, color=GREEN)
    _lab_badge(slide, 1, GREEN)
    _slide_title(slide, "Lab 01 — Expected Result")

    _checklist(slide, [
        "ESP32-CAM เชื่อมต่อ Wi-Fi สำเร็จ",
        "Serial Monitor แสดง IP Address",
        "http://<IP>  เปิดได้ใน Browser",
        "http://<IP>/capture  ส่ง JPEG กลับมา",
        "http://<IP>:81/stream  แสดง Live Video",
    ], top=1.85)
    _tb(slide, "🎉  Lab 01 สำเร็จ! พร้อมสำหรับ Lab 02",
        0.5, 6.72, 12.3, 0.42, size=16, color=GREEN, bold=True, font=FH)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 5 — LAB 02
# ═══════════════════════════════════════════════════════════════════════════════

def s05_header(prs):
    _lab_header(prs, 2, "REST API & Backend",
                "รัน Backend • ทดสอบ API • Mock Provider", "~60 นาที", CYAN)


def s05_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Lab 02 — Overview", "สิ่งที่เราจะทำใน Lab นี้")

    _checklist(slide, [
        "สร้าง Python Virtual Environment (.venv)",
        "ติดตั้ง dependencies ด้วย pip",
        "ตั้งค่า .env ไฟล์  (VISION_PROVIDER=mock)",
        "รัน FastAPI server ด้วย uvicorn",
        "ทดสอบ upload รูปและดูผลผ่าน Web UI",
    ], top=1.85)
    _tb(slide, "🎯  เป้าหมาย: เข้าใจ API data flow — ไม่ต้องใช้ API key",
        0.5, 6.72, 12.3, 0.42, size=14, color=GREEN, bold=True, font=FB)


def s05_step1(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Step 1 — Create Virtual Environment")
    _step_badge(slide, 1, "สร้าง Python Virtual Environment", CYAN)

    code_win = """\
# Windows
cd code\\server
python -m venv .venv
.venv\\Scripts\\activate
# เห็น (.venv) หน้า prompt = สำเร็จ ✅"""
    code_mac = """\
# macOS / Linux
cd code/server
python3 -m venv .venv
source .venv/bin/activate
# เห็น (.venv) หน้า prompt = สำเร็จ ✅"""

    _tb(slide, "Windows:", 0.5, 2.42, 6.3, 0.35, size=14, color=CYAN, bold=True, font=FH)
    _code_box(slide, code_win, 0.5, 2.8, 6.3, 2.55)
    _tb(slide, "macOS / Linux:", 7.0, 2.42, 6.1, 0.35, size=14, color=CYAN, bold=True, font=FH)
    _code_box(slide, code_mac, 7.0, 2.8, 6.1, 2.55)

    _tb(slide, "💡  Virtual Environment แยก Python packages ของ project ออกจากระบบ",
        0.5, 5.52, 12.5, 0.42, size=13, color=GRAY, italic=True, font=FB)


def s05_step2(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Step 2 — Install Dependencies")
    _step_badge(slide, 2, "ติดตั้ง Python Packages", CYAN)

    code = """\
pip install -r requirements.txt

# packages:
#  fastapi          → web framework
#  uvicorn          → ASGI server
#  httpx            → fetch รูปจาก ESP32-CAM
#  python-dotenv    → อ่านค่าจาก .env
#  python-multipart → รับไฟล์ upload
#  pillow           → image processing
#  google-genai     → Gemini Vision SDK
#  anthropic        → Claude Vision SDK"""
    _code_box(slide, code, 0.5, 2.42, 12.3, 4.62)


def s05_step3(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Step 3 — Configure .env")
    _step_badge(slide, 3, "ตั้งค่าไฟล์ Environment", CYAN)

    _tb(slide, "คัดลอก  .env.example  →  .env  แล้วแก้ไข:",
        0.5, 2.42, 12.5, 0.4, size=15, color=LIGHT, font=FB)
    code = """\
# .env  (Lab 02)
VISION_PROVIDER=mock

# เติมทีหลังใน Lab 03:
GEMINI_API_KEY=your_key_here
GEMINI_MODEL=gemini-2.5-flash
CAM_URL=http://192.168.x.x"""
    _code_box(slide, code, 0.5, 2.98, 12.3, 2.8)

    _multiline(slide, [
        ("⚠️  .env อยู่ใน .gitignore — ไม่ถูก commit ขึ้น GitHub", ORANGE, 15, True),
        ("✅  mock provider ไม่ต้องใช้ API key", GREEN, 15, False),
    ], 0.5, 5.95, 12.5, 1.0, gap=10, font=FB)


def s05_step4(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Step 4 — Run FastAPI Server")
    _step_badge(slide, 4, "รัน Backend Server", CYAN)

    code = """\
# ใน terminal (activate .venv แล้ว)
uvicorn main:app --reload --host 0.0.0.0

INFO:     Uvicorn running on http://0.0.0.0:8000
INFO:     Application startup complete.

# เปิดใน Browser:
http://localhost:8000          → Web UI
http://localhost:8000/health   → { "vision_provider": "mock" }
http://localhost:8000/prompts  → preset prompts"""
    _code_box(slide, code, 0.5, 2.42, 12.3, 4.62)


def s05_result(prs):
    slide = _new_slide(prs)
    _top_bar(slide, GREEN); _accent_bar(slide, color=GREEN)
    _lab_badge(slide, 2, GREEN)
    _slide_title(slide, "Lab 02 — Expected Result")

    _checklist(slide, [
        "pip install สำเร็จ ไม่มี error",
        ".env ไฟล์สร้างแล้ว  VISION_PROVIDER=mock",
        "uvicorn รันได้ที่  http://localhost:8000",
        "Web UI เปิดได้และมี preset prompts dropdown",
        "Upload รูปทดสอบ → ได้รับ mock response",
        "/health แสดง  vision_provider: mock",
    ], top=1.85)
    _tb(slide, "🎉  Lab 02 สำเร็จ! พร้อม connect AI จริงใน Lab 03",
        0.5, 6.72, 12.3, 0.42, size=16, color=GREEN, bold=True, font=FH)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 6 — LAB 03
# ═══════════════════════════════════════════════════════════════════════════════

def s06_header(prs):
    _lab_header(prs, 3, "AI Vision Analysis",
                "API Key จริง • ภาพจากกล้อง • Prompt Engineering", "~75 นาที", ORANGE)


def s06_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Lab 03 — Overview", "เชื่อม Vision AI จริง กับกล้อง ESP32-CAM")

    _checklist(slide, [
        "สมัครและรับ Gemini API Key จาก Google AI Studio",
        "ตั้งค่า GEMINI_API_KEY และ VISION_PROVIDER=gemini",
        "Restart uvicorn server",
        "ชี้กล้องและกด Capture & Analyze ครั้งแรก",
        "ทดลอง prompt อย่างน้อย 3 แบบ",
    ], top=1.85)
    _tb(slide, "🎯  เป้าหมาย: AI จริงวิเคราะห์ภาพจากกล้อง ESP32-CAM",
        0.5, 6.72, 12.3, 0.42, size=14, color=GREEN, bold=True, font=FB)


def s06_step1(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Step 1 — Get Gemini API Key")
    _step_badge(slide, 1, "รับ API Key จาก Google AI Studio")

    steps_txt = """\
1. เปิด Browser → aistudio.google.com
2. Sign in ด้วย Google Account
3. คลิก  "Get API key"
4. คลิก  "Create API key"
5. คัดลอก key (ขึ้นต้นด้วย AIza...)
6. วางใน .env:
   GEMINI_API_KEY=AIzaXXXXXXXXXX"""
    _code_box(slide, steps_txt, 0.5, 2.42, 7.5, 4.0)

    _multiline(slide, [
        ("Free Tier:", CYAN, 16, True),
        ("1,500 req / วัน", GREEN, 15, False),
        ("ไม่ต้องใช้บัตรเครดิต", GREEN, 15, False),
        ("", GRAY, 8, False),
        ("⚠️  อย่าแชร์ API key", ORANGE, 15, True),
        ("อย่า commit ลง GitHub", ORANGE, 14, False),
    ], 8.3, 2.42, 4.8, 4.0, gap=9, font=FB)


def s06_step2(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Step 2 — Update .env & Restart")
    _step_badge(slide, 2, "แก้ .env แล้ว Restart Server")

    code = """\
# .env  (Lab 03)
VISION_PROVIDER=gemini
GEMINI_API_KEY=AIzaXXXXXXXXXXXXXXXXXX
GEMINI_MODEL=gemini-2.5-flash
CAM_URL=http://192.168.x.x   ← IP จาก Lab 01

# Restart: Ctrl+C แล้วรันใหม่
uvicorn main:app --reload --host 0.0.0.0

# ตรวจสอบ:
# GET /health → { "vision_provider": "gemini" }"""
    _code_box(slide, code, 0.5, 2.42, 12.3, 4.62)


def s06_step3(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Step 3 — First AI Analysis!")
    _step_badge(slide, 3, "วิเคราะห์ภาพครั้งแรกด้วย Gemini")

    code = """\
1. เปิด http://localhost:8000
2. กรอก ESP32-CAM URL:  http://192.168.x.x
3. เลือก Prompt: "อธิบายสิ่งที่เห็นในภาพนี้"
4. กดปุ่ม  [ Capture & Analyze ]
5. รอ 2-5 วินาที...
6. ดูผลลัพธ์ใน Result panel

# Error codes:
# 401 → API key ผิด / ไม่ได้ตั้ง
# 429 → quota เต็ม → รอ 1 นาที หรือ mock
# 502 → network / camera problem"""
    _code_box(slide, code, 0.5, 2.42, 12.3, 4.62)


def s06_step4(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Step 4 — Prompt Experiments")
    _step_badge(slide, 4, "ทดลอง Prompt อย่างน้อย 3 แบบ")

    prompts = [
        ("Prompt A (กว้าง)",   "อธิบายสิ่งที่เห็นในภาพนี้"),
        ("Prompt B (เจาะจง)",  "นับจำนวนวัตถุแต่ละประเภทในภาพ ตอบเป็น list"),
        ("Prompt C (มี Role)", "คุณคือผู้เชี่ยวชาญ วิเคราะห์ภาพและให้คะแนน 1-10"),
    ]
    top = 2.1
    for label, prompt in prompts:
        _box(slide, 0.5, top, 12.3, 1.3, BG2, border=True)
        _tb(slide, label,         0.75, top + 0.08, 3.5, 0.4, size=15, color=ORANGE, bold=True, font=FH)
        _tb(slide, f'"{prompt}"', 0.75, top + 0.52, 11.8, 0.55, size=15, color=GREEN, font=FB)
        top += 1.5

    _tb(slide, "สังเกต: Prompt ที่ต่างกัน → คำตอบที่ต่างกัน",
        0.5, 6.72, 12.3, 0.42, size=14, color=CYAN, bold=True, font=FB)


def s06_result(prs):
    slide = _new_slide(prs)
    _top_bar(slide, GREEN); _accent_bar(slide, color=GREEN)
    _lab_badge(slide, 3, GREEN)
    _slide_title(slide, "Lab 03 — Expected Result")

    # AI demo image on right
    _photo(slide, "AI Vision Analysis Demo (Lab 03  Cover Result).jpeg",
           7.5, 1.55, 5.5, 5.55, "AI Vision Demo")

    _checklist(slide, [
        "Gemini API Key ตั้งค่าใน .env สำเร็จ",
        "/health แสดง  vision_provider: gemini",
        "Capture & Analyze ได้ผลจาก Gemini จริง",
        "ทดลองอย่างน้อย 3 prompt แตกต่างกัน",
        "Prompt ส่งผลต่อความละเอียดของคำตอบ",
    ], left=0.5, top=1.85)

    _tb(slide, "🎉  Lab 03 สำเร็จ!",
        0.5, 6.72, 7.0, 0.42, size=16, color=GREEN, bold=True, font=FH)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 7 — LAB 04
# ═══════════════════════════════════════════════════════════════════════════════

def s07_header(prs):
    _lab_header(prs, 4, "AI Vision Mini Project",
                "เลือก Project • เขียน Prompt • สร้าง AI ของตัวเอง", "~75 นาที", CYAN)


def s07_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _lab_badge(slide, 4, CYAN)
    _slide_title(slide, "Lab 04 — Overview", "สร้าง AI Vision Project ของตัวเอง")

    flow = [
        ("1. เลือกหัวข้อ",  "เลือก 1 ใน 5 preset หรือคิดขึ้นเอง"),
        ("2. เขียน Prompt", "ใช้สูตร Role + Task + Format + Language"),
        ("3. ทดสอบ",        "กด Capture & Analyze ≥ 3 ครั้ง"),
        ("4. ปรับปรุง",      "แก้ prompt จนได้ผลที่พอใจ"),
        ("5. นำเสนอ",       "demo + อธิบาย prompt + สิ่งที่เรียนรู้ (2 นาที)"),
    ]
    top = 1.75
    for step, desc in flow:
        _list_row(slide, step, desc, 0.5, top, accent=CYAN)
        top += 1.05


def s07_step1(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _lab_badge(slide, 4, CYAN)
    _slide_title(slide, "Step 1 — Choose Your Project")
    _step_badge(slide, 1, "เลือกหัวข้อ Mini Project", CYAN)

    projects = [
        ("🗑️  Waste Classification",  "จำแนกขยะเป็น รีไซเคิล / อินทรีย์ / อันตราย"),
        ("🌿  Plant Health Analysis", "ประเมินสุขภาพต้นไม้จากสี ใบ และลักษณะ"),
        ("🔢  Object Counting",       "นับจำนวนวัตถุแต่ละประเภทในภาพ"),
        ("🖥️  Desk Cleanliness",     "ให้คะแนนความสะอาด / เป็นระเบียบของโต๊ะ"),
        ("🔧  Equipment Detection",   "ตรวจสอบอุปกรณ์ในห้องปฏิบัติการ"),
    ]
    top = 2.42
    for icon_name, desc in projects:
        _tb(slide, icon_name, 0.75, top, 4.5, 0.42, size=16, color=ORANGE, bold=True, font=FH)
        _tb(slide, desc,      5.4,  top, 7.7, 0.42, size=14, color=LIGHT, font=FB)
        top += 0.74

    _tb(slide, "💡  หรือคิดหัวข้อขึ้นเองก็ได้!",
        0.5, 6.75, 12.3, 0.42, size=14, color=CYAN, italic=True, font=FB)


def s07_step2(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _lab_badge(slide, 4, CYAN)
    _slide_title(slide, "Step 2 — Write Your Prompt")
    _step_badge(slide, 2, "เขียน Custom Prompt สำหรับ Project", CYAN)

    code = """\
# สูตร: Role + Task + Format + Language

คุณคือ [บทบาทของ AI]

วิเคราะห์ภาพนี้และ[สิ่งที่ต้องทำ]

ตอบในรูปแบบ:
[หัวข้อ 1]: [ค่า]
[หัวข้อ 2]: [ค่า]
สรุป: [ประโยคสรุป]

ตอบเป็นภาษาไทย"""
    _code_box(slide, code, 0.5, 2.45, 7.5, 4.5)

    _multiline(slide, [
        ("Tips:", CYAN, 16, True),
        ("• ทดสอบอย่างน้อย 3 ครั้ง", LIGHT, 14, False),
        ("• AI ตอบกว้างเกิน", LIGHT, 14, False),
        ("  → เพิ่มรายละเอียดใน Format", GRAY, 13, False),
        ("• AI ตอบไม่ตรง", LIGHT, 14, False),
        ("  → ปรับ Role ให้ชัดขึ้น", GRAY, 13, False),
    ], 8.2, 2.45, 4.9, 4.5, gap=9, font=FB)


def s07_tips(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _lab_badge(slide, 4, CYAN)
    _slide_title(slide, "Presentation Tips", "วิธีนำเสนอ Mini Project (2 นาที)")

    items = [
        ("🎯  Demo Live",     "ชี้กล้องไปที่วัตถุจริง แล้วกด Capture & Analyze"),
        ("📝  แสดง Prompt",   "บอกว่าเขียน prompt อะไร และทำไมถึงเลือกแบบนั้น"),
        ("📊  แสดงผลลัพธ์",  "AI ตอบอะไร? ตรงกับที่คาดหวังไหม?"),
        ("💡  สิ่งที่เรียนรู้", "ได้เรียนรู้อะไรจาก workshop นี้?"),
    ]
    top = 1.88
    for title, desc in items:
        _list_row(slide, title, desc, 0.5, top, accent=CYAN)
        top += 1.25


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 8 — MINI PROJECTS
# ═══════════════════════════════════════════════════════════════════════════════

def s08_header(prs):
    _section_header(prs, 8, "Mini Projects", "ไอเดียโปรเจกต์สำหรับ Lab 04", ORANGE)


def _project_slide(prs, icon, title_en, title_th, prompt, output):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE); _accent_bar(slide, color=ORANGE)
    _tb(slide, f"{icon}  {title_en}", 0.5, 0.2, 12.5, 0.7,
        size=26, color=ORANGE, bold=True, font=FH)
    _tb(slide, title_th, 0.5, 0.88, 12.5, 0.42,
        size=15, color=LIGHT, font=FB)
    _box(slide, 0.5, 1.42, 12.3, 0.04, ORANGE)

    _tb(slide, "Prompt:", 0.5, 1.58, 2.0, 0.38, size=14, color=CYAN, bold=True, font=FH)
    _code_box(slide, prompt, 0.5, 1.98, 12.3, 2.52)

    _tb(slide, "Expected Output:", 0.5, 4.62, 3.5, 0.38, size=14, color=GREEN, bold=True, font=FH)
    _code_box(slide, output, 0.5, 5.02, 12.3, 2.12)


def s08_waste(prs):
    _project_slide(prs,
        "🗑️", "Waste Classification", "จำแนกประเภทขยะ",
        'คุณคือระบบ AI จำแนกขยะอัจฉริยะ\nวิเคราะห์ขยะในภาพและตอบ:\n🗑️ ประเภทขยะ: [ชื่อ]\n♻️ หมวดหมู่: รีไซเคิล / อินทรีย์ / ทั่วไป / อันตราย\n✅ วิธีทิ้ง: [คำแนะนำ]\nตอบภาษาไทย',
        '🗑️ ประเภทขยะ: ขวดพลาสติก\n♻️ หมวดหมู่: รีไซเคิล\n✅ วิธีทิ้ง: ทิ้งในถังขยะสีเหลือง ล้างให้สะอาดก่อนทิ้ง')


def s08_plant(prs):
    _project_slide(prs,
        "🌿", "Plant Health Analysis", "วิเคราะห์สุขภาพต้นไม้",
        'คุณคือนักพฤกษศาสตร์ผู้เชี่ยวชาญ\nวิเคราะห์ต้นไม้ในภาพและตอบ:\n🌿 ชนิด: [ถ้าระบุได้]\n💚 สุขภาพ: ดีมาก / ดี / ปานกลาง / ต้องการดูแลด่วน\n🔍 สังเกต: [สีใบ, จุด, ใบเหี่ยว]\n💧 การรดน้ำ: [แนะนำ]\n🌱 คำแนะนำ: [2-3 ข้อ]',
        '🌿 ชนิด: ต้นมอนสเตอร่า\n💚 สุขภาพ: ดี\n🔍 สังเกต: ใบสีเขียวเข้ม สมบูรณ์\n🌱 คำแนะนำ: วางในที่มีแสงรำไร')


def s08_count(prs):
    _project_slide(prs,
        "🔢", "Object Counting", "นับจำนวนวัตถุ",
        'นับและจำแนกวัตถุทุกชิ้นในภาพอย่างละเอียด\nตอบในรูปแบบ:\n[ชื่อวัตถุ]: [จำนวน]\nรวมทั้งหมด: [จำนวน] ชิ้น\nหมายเหตุ: [วัตถุที่มองไม่ชัด]',
        'โน้ตบุ๊ค: 1\nกาแฟ: 1\nหนังสือ: 3\nปากกา: 2\nรวมทั้งหมด: 7 ชิ้น')


def s08_desk(prs):
    _project_slide(prs,
        "🖥️", "Desk Cleanliness Inspection", "ประเมินความเป็นระเบียบของโต๊ะ",
        'คุณคือผู้ตรวจสอบความสะอาด ประเมินโต๊ะนี้:\n⭐ คะแนน: [1-5 ดาว]\n📦 สิ่งของ: [รายการ]\n✅ จุดที่ดี: [สิ่งที่เป็นระเบียบ]\n⚠️ ปรับปรุง: [สิ่งที่ควรจัดใหม่]',
        '⭐ คะแนน: 4 ดาว\n📦 สิ่งของ: โน้ตบุ๊ค, เมาส์, แก้วน้ำ\n✅ จุดที่ดี: จัดสายไฟเป็นระเบียบ\n⚠️ ปรับปรุง: มีเศษกระดาษบนโต๊ะ')


def s08_equipment(prs):
    _project_slide(prs,
        "🔧", "Equipment Detection", "ตรวจสอบอุปกรณ์",
        'คุณคือระบบตรวจสอบอุปกรณ์ห้องปฏิบัติการ\nวิเคราะห์ภาพและระบุ:\n🔧 อุปกรณ์ที่พบ: [รายการ]\n✅ พร้อมใช้งาน: [รายการ]\n⚠️ ต้องตรวจสอบ: [รายการ]\nสรุป: [1 ประโยค]',
        '🔧 พบ: ESP32-CAM, สาย USB, breadboard\n✅ พร้อม: ESP32-CAM, สาย USB\n⚠️ ตรวจสอบ: breadboard (มีออกซิไดซ์)\nสรุป: อุปกรณ์ส่วนใหญ่พร้อมใช้งาน')


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 9 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════

def s09_careers(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Future Career Paths", "เส้นทางอาชีพที่เกี่ยวข้อง")

    careers = [
        ("🤖  AI Engineer",          "พัฒนาโมเดล AI, fine-tuning, deployment"),
        ("📡  IoT Developer",         "ออกแบบระบบเซนเซอร์, firmware, cloud"),
        ("📊  Data Scientist",        "วิเคราะห์ข้อมูล, ML models, visualization"),
        ("🤖  Robotics Engineer",     "หุ่นยนต์ที่ใช้ computer vision นำทาง"),
        ("🛡️  AI Safety Researcher", "ความปลอดภัยและจริยธรรมของ AI"),
        ("🚀  AI Product Manager",    "นำทีมสร้าง product ที่ขับเคลื่อนด้วย AI"),
    ]
    top = 1.72
    for i, (title, desc) in enumerate(careers):
        col = 0 if i % 2 == 0 else 6.6
        row = top + (i // 2) * 1.62
        _box(slide, 0.5 + col, row, 6.3, 1.38, BG2, border=True)
        _tb(slide, title, 0.75 + col, row + 0.1,  5.8, 0.45, size=16, color=CYAN, bold=True, font=FH)
        _tb(slide, desc,  0.75 + col, row + 0.58, 5.8, 0.55, size=13, color=LIGHT, font=FB)


def s09_resources(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Additional Resources", "แหล่งเรียนรู้เพิ่มเติม")

    resources = [
        ("📖  Workshop GitHub",        "github.com/Pachara966/esp32cam-ai-vision-workshop"),
        ("🤖  Google AI Studio",       "aistudio.google.com  (Gemini API Key)"),
        ("🧠  Anthropic Console",      "console.anthropic.com  (Claude API Key)"),
        ("📚  FastAPI Docs",           "fastapi.tiangolo.com"),
        ("📷  ESP32 Arduino Docs",     "docs.espressif.com"),
        ("🎓  Google ML Crash Course", "developers.google.com/machine-learning/crash-course"),
    ]
    top = 1.72
    for title, url in resources:
        _list_row(slide, title, url, 0.5, top, accent=ORANGE, title_color=ORANGE)
        top += 0.95


def s09_summary(prs):
    slide = _new_slide(prs)
    _top_bar(slide); _accent_bar(slide)
    _slide_title(slide, "Workshop Summary", "สรุป Workshop วันนี้")

    summary = [
        ("Lab 01", "ติดตั้ง ESP32-CAM + Flash Firmware + Live Stream",     ORANGE),
        ("Lab 02", "รัน FastAPI Backend + ทำความเข้าใจ REST API",            CYAN),
        ("Lab 03", "เชื่อม Gemini Vision + Prompt Engineering",             ORANGE),
        ("Lab 04", "สร้าง AI Vision Mini Project ของตัวเอง",                GREEN),
    ]
    top = 1.82
    for lab, desc, color in summary:
        _box(slide, 0.5, top, 12.3, 1.05, BG2, border=True)
        _box(slide, 0.5, top, 1.5, 1.05, color)
        _tb(slide, lab, 0.5, top + 0.2, 1.5, 0.55,
            size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FH)
        _tb(slide, desc, 2.2, top + 0.22, 10.5, 0.55,
            size=16, color=LIGHT, font=FB)
        top += 1.22

    _tb(slide, "🎉  ภายใน 1 วัน นักเรียนสร้างระบบ AI Vision ที่ใช้งานได้จริง!",
        0.5, 6.62, 12.3, 0.5, size=16, color=GREEN, bold=True,
        align=PP_ALIGN.CENTER, font=FH)


def s09_qa(prs):
    slide = _new_slide(prs)
    _box(slide, 0, 0, 13.33, 7.5, BG)
    _box(slide, 0, 0,    13.33, 0.09, CYAN)
    _box(slide, 0, 7.41, 13.33, 0.09, CYAN)

    _tb(slide, "?", 2.5, 0.5, 8.33, 5.8,
        size=260, color=GHOST, bold=True, align=PP_ALIGN.CENTER, font=FH)
    _tb(slide, "Q&A", 0.5, 2.0, 12.33, 2.5,
        size=100, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FH)
    _tb(slide, "ถามได้เลยครับ / Feel free to ask!",
        0.5, 5.8, 12.33, 0.6, size=22, color=CYAN, align=PP_ALIGN.CENTER, font=FB)


def s09_thankyou(prs):
    slide = _new_slide(prs)

    # Group photo right panel
    _photo(slide, "Workshop Group Photo Style (Thank You Slide).jpeg",
           6.5, 0, 6.83, 7.5)
    _box(slide, 6.2, 0, 0.6, 7.5, BG)  # fade edge

    _box(slide, 0, 0,    13.33, 0.12, ORANGE)
    _box(slide, 0, 7.38, 13.33, 0.12, ORANGE)
    _box(slide, 0, 0.12, 0.1,   7.26, CYAN)

    _tb(slide, "ขอบคุณทุกคน!", 0.5, 1.0, 6.0, 1.2,
        size=54, color=WHITE, bold=True, font=FH)
    _tb(slide, "Thank You!", 0.5, 2.1, 6.0, 0.8,
        size=38, color=CYAN, bold=True, font=FH)

    _box(slide, 0.5, 3.05, 4.5, 0.05, GRAY)

    _tb(slide, "มหาวิทยาลัยเทคโนโลยีราชมงคลธัญบุรี (RMUTT)",
        0.5, 3.22, 6.0, 0.5, size=17, color=LIGHT, font=FB)
    _tb(slide, "[ชื่อผู้สอน / Instructor Name]",
        0.5, 3.78, 6.0, 0.4, size=14, color=GRAY, font=FB)
    _tb(slide, "github.com/Pachara966/esp32cam-ai-vision-workshop",
        0.5, 4.62, 6.0, 0.42, size=13, color=CYAN, font=FB)
    _tb(slide, "#ESP32CAM_AI_RMUTT",
        0.5, 5.12, 6.0, 0.42, size=17, color=ORANGE, bold=True, font=FH)
    _tb(slide, "ESP32-CAM AI Vision Workshop Team",
        0.5, 6.55, 6.0, 0.42, size=12, color=GRAY, font=FB)


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Building Workshop_v3 (dark navy + Sarabun + modern layout + images)...")

    s00_cover(prs);     print("  01 Cover")
    s00_overview(prs);  print("  02 Overview")
    s00_outcomes(prs);  print("  03 Outcomes")
    s00_agenda(prs);    print("  04 Agenda")

    s01_header(prs);   print("  05 Sec 1 Header")
    s01_iot(prs);      print("  06 IoT")
    s01_esp32cam(prs); print("  07 ESP32-CAM")
    s01_hardware(prs); print("  08 Hardware")
    s01_wifi(prs);     print("  09 Wi-Fi")

    s02_header(prs);   print("  10 Sec 2 Header")
    s02_api(prs);      print("  11 REST API")
    s02_http(prs);     print("  12 HTTP")
    s02_json(prs);     print("  13 JSON")
    s02_fastapi(prs);  print("  14 FastAPI")

    s03_header(prs);   print("  15 Sec 3 Header")
    s03_ai(prs);       print("  16 AI")
    s03_genai(prs);    print("  17 GenAI")
    s03_vision(prs);   print("  18 Vision")
    s03_gemini(prs);   print("  19 Gemini")
    s03_claude(prs);   print("  20 Claude")
    s03_prompt(prs);   print("  21 Prompt")
    s03_workflow(prs); print("  22 Workflow")

    s04_header(prs);   print("  23 Lab01 Header")
    s04_overview(prs); print("  24 Lab01 Overview")
    s04_step1(prs);    print("  25 Lab01 Step1")
    s04_step2(prs);    print("  26 Lab01 Step2")
    s04_step3(prs);    print("  27 Lab01 Step3")
    s04_step4(prs);    print("  28 Lab01 Step4")
    s04_step5(prs);    print("  29 Lab01 Step5")
    s04_result(prs);   print("  30 Lab01 Result")

    s05_header(prs);   print("  31 Lab02 Header")
    s05_overview(prs); print("  32 Lab02 Overview")
    s05_step1(prs);    print("  33 Lab02 Step1")
    s05_step2(prs);    print("  34 Lab02 Step2")
    s05_step3(prs);    print("  35 Lab02 Step3")
    s05_step4(prs);    print("  36 Lab02 Step4")
    s05_result(prs);   print("  37 Lab02 Result")

    s06_header(prs);   print("  38 Lab03 Header")
    s06_overview(prs); print("  39 Lab03 Overview")
    s06_step1(prs);    print("  40 Lab03 Step1")
    s06_step2(prs);    print("  41 Lab03 Step2")
    s06_step3(prs);    print("  42 Lab03 Step3")
    s06_step4(prs);    print("  43 Lab03 Step4")
    s06_result(prs);   print("  44 Lab03 Result")

    s07_header(prs);   print("  45 Lab04 Header")
    s07_overview(prs); print("  46 Lab04 Overview")
    s07_step1(prs);    print("  47 Lab04 Step1")
    s07_step2(prs);    print("  48 Lab04 Step2")
    s07_tips(prs);     print("  49 Lab04 Tips")

    s08_header(prs);    print("  50 MiniProj Header")
    s08_waste(prs);     print("  51 Waste")
    s08_plant(prs);     print("  52 Plant")
    s08_count(prs);     print("  53 Count")
    s08_desk(prs);      print("  54 Desk")
    s08_equipment(prs); print("  55 Equipment")

    s09_careers(prs);   print("  56 Careers")
    s09_resources(prs); print("  57 Resources")
    s09_summary(prs);   print("  58 Summary")
    s09_qa(prs);        print("  59 Q&A")
    s09_thankyou(prs);  print("  60 Thank You")

    out = "Workshop_v3.pptx"
    prs.save(out)
    print(f"\nSaved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
