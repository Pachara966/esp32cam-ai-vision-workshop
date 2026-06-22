#!/usr/bin/env python3
"""
Generate Workshop.pptx — ESP32-CAM AI Vision Workshop (RMUTT)
Usage : python generate_pptx.py   (run from inside slides/ directory)
Output: Workshop.pptx  (~60 slides, dark-navy tech theme, bilingual TH/EN)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy – slide background
BG2     = RGBColor(0x0A, 0x13, 0x1E)   # deeper navy – code panels, ghost
GHOST   = RGBColor(0x17, 0x2A, 0x3E)   # section-number ghost
CYAN    = RGBColor(0x00, 0xC8, 0xFF)   # primary accent
ORANGE  = RGBColor(0xFF, 0x70, 0x43)   # lab / callout accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xD0, 0xE4, 0xF5)   # body text
GRAY    = RGBColor(0x88, 0xA0, 0xB8)   # secondary / italic text
GREEN   = RGBColor(0x4C, 0xD9, 0x64)   # success / checkmarks

# ── Slide dimensions (16 : 9) ─────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
#  LOW-LEVEL HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def _tb(slide, text, left, top, width, height,
        size=18, color=LIGHT, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True):
    """Single-paragraph text box."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def _box(slide, left, top, width, height, fill_color, border=False):
    """Filled rectangle (no border by default)."""
    s = slide.shapes.add_shape(
        1,   # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border:
        s.line.color.rgb = CYAN
    else:
        s.line.fill.background()
    return s


def _multiline(slide, lines, left, top, width, height,
               size=18, color=LIGHT, bold=False,
               align=PP_ALIGN.LEFT, gap=6):
    """
    Multi-paragraph text box.
    lines = list of:
      str                → plain line in `color`
      (str, RGBColor)    → line with specific colour
      (str, RGBColor, int, bool)  → text, colour, size, bold
    """
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, col, sz, bd = item, color, size, bold
        elif len(item) == 2:
            text, col = item; sz, bd = size, bold
        elif len(item) == 3:
            text, col, sz = item; bd = bold
        else:
            text, col, sz, bd = item

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.bold = bd
    return tb


def _accent_bar(slide, y=7.25, color=CYAN):
    """Thin horizontal accent bar."""
    _box(slide, 0, y, 13.33, 0.08, color)


def _top_bar(slide, color=CYAN):
    _box(slide, 0, 0, 13.33, 0.07, color)


def _section_ghost(slide, num):
    """Large ghost section number — decorative."""
    _tb(slide, str(num), 0.1, 0.5, 4.5, 6.5,
        size=220, color=GHOST, bold=True, align=PP_ALIGN.LEFT)


def _lab_badge(slide, num, color=ORANGE):
    """Coloured lab badge – top right corner."""
    _box(slide, 11.0, 0.18, 2.1, 0.62, color)
    _tb(slide, f"LAB  {num:02d}", 11.0, 0.18, 2.1, 0.62,
        size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def _slide_title(slide, text, sub=None):
    """Standard slide title + optional sub-label."""
    _tb(slide, text, 0.5, 0.35, 12.3, 0.9,
        size=36, color=WHITE, bold=True)
    _box(slide, 0.5, 1.18, 3.0, 0.06, CYAN)   # underline
    if sub:
        _tb(slide, sub, 0.5, 1.3, 12.3, 0.5,
            size=16, color=GRAY, italic=True)


def _code_box(slide, code, left, top, width, height):
    """Dark monospace code panel."""
    _box(slide, left, top, width, height, BG2, border=True)
    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12),
        Inches(width - 0.3), Inches(height - 0.25))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.size = Pt(14)
        r.font.color.rgb = GREEN
        r.font.name = "Courier New"


def _step_badge(slide, step_num, label="", color=ORANGE):
    """Coloured step circle + label."""
    _box(slide, 0.5, 1.6, 0.7, 0.7, color)
    _tb(slide, str(step_num), 0.5, 1.6, 0.7, 0.7,
        size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if label:
        _tb(slide, label, 1.35, 1.65, 11.5, 0.6,
            size=22, color=WHITE, bold=True)


def _checklist(slide, items, left=0.8, top=2.0, size=20):
    """Green checkmark bullet list."""
    lines = []
    for item in items:
        lines.append((f"✅  {item}", WHITE, size, False))
    _multiline(slide, lines, left, top, 12.0, 5.0, gap=14)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 0  —  OPENING
# ═══════════════════════════════════════════════════════════════════════════════

def s00_cover(prs):
    slide = _new_slide(prs)
    # Top / bottom accent bars
    _box(slide, 0, 0, 13.33, 0.12, CYAN)
    _box(slide, 0, 7.38, 13.33, 0.12, CYAN)
    # Decorative side stripe
    _box(slide, 0, 0.12, 0.08, 7.26, ORANGE)

    _tb(slide, "ESP32-CAM AI Vision Workshop",
        1.0, 1.0, 11.33, 1.2,
        size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _tb(slide, "From IoT Camera to Intelligent Image Analysis",
        1.0, 2.1, 11.33, 0.7,
        size=22, color=CYAN, align=PP_ALIGN.CENTER)

    _box(slide, 3.5, 2.95, 6.33, 0.05, GRAY)   # divider

    _tb(slide, "การพัฒนากล้องอัจฉริยะด้วย ESP32-CAM และปัญญาประดิษฐ์",
        1.0, 3.1, 11.33, 0.7,
        size=20, color=LIGHT, align=PP_ALIGN.CENTER)

    _tb(slide, "สำหรับการวิเคราะห์ภาพ",
        1.0, 3.7, 11.33, 0.5,
        size=20, color=LIGHT, align=PP_ALIGN.CENTER)

    _tb(slide, "มหาวิทยาลัยเทคโนโลยีราชมงคลธัญบุรี  (RMUTT)",
        1.0, 4.6, 11.33, 0.55,
        size=18, color=GRAY, align=PP_ALIGN.CENTER)

    _tb(slide, "[วันที่ / Date]     •     #ESP32CAM_AI_RMUTT",
        1.0, 5.2, 11.33, 0.45,
        size=15, color=GRAY, align=PP_ALIGN.CENTER)


def s00_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Workshop Overview", "ในวันนี้เราจะสร้างอะไร?")

    _multiline(slide, [
        ("เราจะสร้างระบบกล้อง AI ที่สามารถ...", CYAN, 20, True),
        ("  วิเคราะห์ภาพ  •  นับวัตถุ  •  จำแนกประเภท  •  ตอบคำถาม", LIGHT, 18, False),
    ], 0.5, 1.85, 12.5, 0.9, gap=4)

    # Flow diagram
    _box(slide, 0.5, 2.95, 12.3, 3.8, BG2, border=True)
    flow = [
        ("📷  ESP32-CAM",        CYAN,   20, True),
        ("        │",            GRAY,   16, False),
        ("        ▼",            GRAY,   16, False),
        ("📶  Wi-Fi  →  /capture endpoint",  LIGHT, 18, False),
        ("        │",            GRAY,   16, False),
        ("        ▼",            GRAY,   16, False),
        ("⚙️   FastAPI Backend  (Python)",    LIGHT, 18, False),
        ("        │",            GRAY,   16, False),
        ("        ▼",            GRAY,   16, False),
        ("🤖  Vision AI  (Gemini / Claude)",  ORANGE, 18, True),
        ("        │",            GRAY,   16, False),
        ("        ▼",            GRAY,   16, False),
        ("📊  Web UI  —  ผลการวิเคราะห์ภาพ",  GREEN,  18, True),
    ]
    _multiline(slide, flow, 3.5, 3.05, 7.0, 3.6, gap=0)


def s00_outcomes(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Learning Outcomes", "เมื่อสิ้นสุด Workshop นักเรียนจะสามารถ:")

    outcomes = [
        ("1.  ติดตั้งและใช้งาน ESP32-CAM",             "Configure & Flash Firmware"),
        ("2.  อธิบาย REST API",                          "Explain REST API & HTTP"),
        ("3.  รัน Python Backend ด้วย FastAPI",          "Run FastAPI Server"),
        ("4.  ใช้ Vision AI API วิเคราะห์ภาพ",           "Use Gemini / Claude Vision"),
        ("5.  เขียน Prompt ที่เจาะจง",                   "Design Effective Prompts"),
        ("6.  สร้าง AI Vision Mini Project",              "Build a Real AI Project"),
    ]
    top = 1.7
    for th, en in outcomes:
        _tb(slide, th, 0.7, top, 12.0, 0.42,
            size=20, color=WHITE, bold=False)
        _tb(slide, f"     {en}", 0.7, top + 0.38, 12.0, 0.3,
            size=13, color=GRAY, italic=True)
        top += 0.82


def s00_agenda(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Workshop Agenda", "ตารางเวลา 1 วัน")

    rows = [
        ("08:30 – 09:00", "ลงทะเบียน / รับอุปกรณ์",            "—"),
        ("09:00 – 09:30", "Lecture: AI + IoT + Computer Vision",  "—"),
        ("09:30 – 10:30", "Lab 01 — ESP32-CAM Camera Server",     "🔌"),
        ("10:30 – 10:45", "พัก ☕",                                "—"),
        ("10:45 – 11:45", "Lab 02 — REST API & Backend",          "⚙️"),
        ("11:45 – 12:45", "พักกลางวัน 🍱",                        "—"),
        ("12:45 – 13:00", "Lecture: Prompt Engineering",           "—"),
        ("13:00 – 14:15", "Lab 03 — AI Vision Analysis",          "🤖"),
        ("14:15 – 14:30", "พัก",                                  "—"),
        ("14:30 – 15:45", "Lab 04 — Mini Project",                 "🚀"),
        ("15:45 – 16:00", "นำเสนอผลงาน / มอบใบประกาศ",           "🎉"),
    ]

    top = 1.55
    for time, activity, icon in rows:
        col = CYAN if "Lab" in activity or "Lecture" in activity else LIGHT
        bd  = "Lab" in activity or "Lecture" in activity
        _tb(slide, time,     0.4, top, 2.3, 0.38, size=14, color=GRAY)
        _tb(slide, activity, 2.8, top, 9.5, 0.38, size=14, color=col, bold=bd)
        top += 0.38


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 1  —  IoT & ESP32-CAM
# ═══════════════════════════════════════════════════════════════════════════════

def _section_header(prs, num, title_en, title_th, accent=CYAN):
    slide = _new_slide(prs)
    _section_ghost(slide, num)
    _box(slide, 0, 0, 0.12, 7.5, accent)           # left stripe
    _box(slide, 0.12, 6.8, 13.21, 0.06, accent)    # bottom stripe
    _tb(slide, title_en, 4.5, 2.2, 8.5, 1.2,
        size=52, color=WHITE, bold=True)
    _tb(slide, title_th,  4.5, 3.4, 8.5, 0.7,
        size=22, color=accent)
    _tb(slide, f"Section {num}", 4.5, 4.2, 8.5, 0.4,
        size=15, color=GRAY)


def s01_header(prs):
    _section_header(prs, 1, "IoT & ESP32-CAM", "อินเทอร์เน็ตของสรรพสิ่ง และกล้อง AI")


def s01_iot(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "What is IoT?", "Internet of Things — อินเทอร์เน็ตของสรรพสิ่ง")

    _tb(slide,
        "IoT = อุปกรณ์ที่เชื่อมต่ออินเทอร์เน็ตเพื่อส่ง/รับข้อมูล",
        0.5, 1.8, 12.5, 0.55, size=22, color=CYAN, bold=True)

    examples = [
        ("🏠  Smart Home",       "หลอดไฟ, กล้องวงจรปิด, ล็อคประตู ที่ควบคุมผ่านมือถือ"),
        ("🏭  Smart Factory",    "เซนเซอร์ตรวจจับอุณหภูมิ, แรงสั่นสะเทือน, แจ้งเตือนอัตโนมัติ"),
        ("🌱  Smart Agriculture","กล้องถ่ายภาพต้นไม้ → AI วิเคราะห์โรค → แจ้งเตือนเกษตรกร"),
    ]
    top = 2.55
    for icon_title, desc in examples:
        _box(slide, 0.5, top, 12.3, 1.1, BG2, border=True)
        _tb(slide, icon_title, 0.75, top + 0.08, 4.0, 0.45,
            size=19, color=ORANGE, bold=True)
        _tb(slide, desc, 0.75, top + 0.52, 11.8, 0.45,
            size=15, color=LIGHT)
        top += 1.28


def s01_esp32cam(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "ESP32-CAM Introduction", "โมดูลกล้อง AI Thinker — หัวใจของ Workshop")

    specs = [
        ("CPU",         "ESP32 Dual-Core 240 MHz"),
        ("RAM",         "520 KB SRAM  +  4 MB PSRAM"),
        ("Flash",       "4 MB"),
        ("Camera",      "OV2640  2MP  (สูงสุด 1600×1200)"),
        ("Wi-Fi",       "802.11 b/g/n  2.4 GHz"),
        ("Price",       "~150 – 250 บาท"),
    ]
    _box(slide, 0.5, 1.7, 7.5, 5.3, BG2, border=True)
    top = 1.9
    for label, val in specs:
        _tb(slide, label, 0.75, top, 2.3, 0.5,
            size=16, color=CYAN, bold=True)
        _tb(slide, val,   3.1,  top, 4.8, 0.5,
            size=16, color=LIGHT)
        top += 0.72

    _tb(slide, "AI Thinker\nESP32-CAM", 8.5, 2.5, 4.5, 3.0,
        size=28, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, "[ รูป / Photo ]", 8.5, 4.0, 4.5, 1.0,
        size=14, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


def s01_hardware(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "ESP32-CAM Hardware Tour", "รู้จักส่วนประกอบสำคัญ")

    diagram = """\
  ┌─────────────────────────────┐
  │        OV2640 Camera        │  ← ถอด-เสียบได้ (อย่าดึงแรง)
  │  ┌───────────────────────┐  │
  │  │    ESP32-WROVER SoC   │  │  ← CPU + Wi-Fi + RAM
  │  └───────────────────────┘  │
  │  GPIO0 ── ต้องต่อ GND ขณะ flash │  ← ⚠️ สำคัญมาก
  │  5V / GND ── จ่ายไฟ          │
  │  GPIO1/3  ── UART (Serial)   │
  └─────────────────────────────┘
         │  USB (ผ่าน MB Board)
         ▼
     ESP32-CAM-MB  (CH340 USB-to-Serial)"""

    _code_box(slide, diagram, 0.5, 1.7, 12.3, 5.3)


def s01_wifi(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Wi-Fi & Networking Basics", "เครือข่ายในห้อง Workshop")

    points = [
        ("SSID",        "ชื่อ Wi-Fi ที่มองเห็น  → ใส่ใน secrets.h"),
        ("Password",    "รหัส Wi-Fi             → ใส่ใน secrets.h"),
        ("IP Address",  "เช่น 192.168.1.50 — ที่อยู่ของ ESP32-CAM ในเครือข่าย"),
        ("Port",        "Web UI: port 80  •  Stream: port 81"),
        ("⚠️  2.4 GHz",  "ESP32 รองรับแค่ 2.4 GHz  —  ไม่รองรับ 5 GHz"),
    ]
    top = 1.8
    for label, desc in points:
        _box(slide, 0.5, top, 12.3, 0.82, BG2)
        _tb(slide, label, 0.75, top + 0.08, 2.8, 0.45,
            size=17, color=CYAN, bold=True)
        _tb(slide, desc,  3.65, top + 0.08, 9.0, 0.6,
            size=16, color=LIGHT)
        top += 0.97


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 2  —  REST API & BACKEND
# ═══════════════════════════════════════════════════════════════════════════════

def s02_header(prs):
    _section_header(prs, 2, "REST API & Backend", "การสื่อสารระหว่างอุปกรณ์กับ AI")


def s02_api(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "What is REST API?", "API คืออะไร? — เปรียบเหมือนร้านอาหาร")

    cols = [
        ("🍽️  ร้านอาหาร",  ["ลูกค้า = Client (เบราว์เซอร์)", "พนักงาน = API (ตัวกลาง)", "ครัว = Server (Backend)", "เมนู = Endpoints", "ออเดอร์ = Request", "อาหาร = Response"]),
        ("💻  REST API",   ["Browser = Client", "API = Interface Layer", "FastAPI = Server", "GET /prompts, POST /analyze", "HTTP Request + JSON Body", "JSON Response"]),
    ]
    for i, (heading, items) in enumerate(cols):
        left = 0.5 + i * 6.6
        _box(slide, left, 1.75, 6.3, 5.3, BG2, border=True)
        _tb(slide, heading, left + 0.2, 1.85, 5.9, 0.5,
            size=20, color=CYAN, bold=True)
        _box(slide, left, 2.3, 6.3, 0.04, CYAN)
        lines = [(item, LIGHT, 16, False) for item in items]
        _multiline(slide, lines, left + 0.2, 2.45, 5.9, 4.4, gap=10)


def s02_http(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "HTTP GET vs POST", "วิธีส่งข้อมูลผ่าน HTTP")

    for i, (method, color, use, example) in enumerate([
        ("GET",  CYAN,   "ขอดึงข้อมูล (Read)",    "GET /prompts\nGET /health"),
        ("POST", ORANGE, "ส่งข้อมูลไป (Write/Send)", "POST /analyze\n{ cam_url, prompt }"),
    ]):
        left = 0.5 + i * 6.6
        _box(slide, left, 1.75, 6.3, 5.4, BG2, border=True)
        _box(slide, left, 1.75, 6.3, 0.62, color)
        _tb(slide, method, left + 0.2, 1.78, 5.9, 0.55,
            size=28, color=WHITE, bold=True)
        _tb(slide, use, left + 0.2, 2.55, 5.9, 0.5,
            size=18, color=LIGHT, bold=True)
        _code_box(slide, example, left + 0.2, 3.2, 5.9, 1.5)
        _tb(slide, "ไม่มี body / URL parameters",
            left + 0.2, 4.85, 5.9, 0.4,
            size=13, color=GRAY, italic=(i == 0))
        _tb(slide, "ส่ง body เป็น JSON" if i == 1 else "",
            left + 0.2, 4.85, 5.9, 0.4,
            size=13, color=GRAY, italic=True)


def s02_json(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "JSON Basics", "รูปแบบข้อมูลที่ใช้สื่อสารระหว่าง Backend กับ AI")

    _tb(slide, "JSON = JavaScript Object Notation — ข้อมูลในรูปแบบ key: value",
        0.5, 1.75, 12.5, 0.5, size=17, color=CYAN)

    sample = """\
{
  "provider": "gemini",
  "filename":  "photo.jpg",
  "image_size_bytes": 12480,
  "result": "ภาพแสดงโต๊ะทำงานที่มีโน้ตบุ๊ค 1 เครื่อง
             กาแฟ 1 แก้ว และหนังสือ 3 เล่ม
             โต๊ะมีความเป็นระเบียบระดับดี"
}"""
    _code_box(slide, sample, 0.5, 2.35, 7.5, 4.7)

    _multiline(slide, [
        ("Key Points:", CYAN, 17, True),
        ("• ใช้  { }  สำหรับ object", LIGHT, 16, False),
        ("• ใช้  [ ]  สำหรับ list/array", LIGHT, 16, False),
        ("• Key ใส่ใน  \" \"  เสมอ", LIGHT, 16, False),
        ("• Value อาจเป็น string, number,", LIGHT, 16, False),
        ("  boolean, list หรือ object", LIGHT, 16, False),
        ("• ทุก API ใน Workshop นี้", LIGHT, 16, False),
        ("  ตอบกลับเป็น JSON", LIGHT, 16, False),
    ], 8.3, 2.35, 4.8, 4.7, gap=8)


def s02_fastapi(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "FastAPI Backend Overview", "สถาปัตยกรรมของ Backend")

    arch = """\
Browser / Web UI  (localhost:8000)
        │
        │  POST /analyze   { cam_url, prompt }
        ▼
  ┌─────────────────────────────────────┐
  │           FastAPI (main.py)         │
  │   GET /       → index.html          │
  │   GET /prompts → preset list        │
  │   POST /analyze → fetch + AI        │
  │   POST /analyze-upload → AI         │
  │   GET /health  → status             │
  └──────────────┬──────────────────────┘
                 │                  │
        ┌────────┴────────┐  ┌──────┴──────┐
        │  ESP32-CAM       │  │  Vision AI  │
        │  /capture (JPEG) │  │  Provider   │
        │  http://IP/...   │  │  mock/gemini│
        └──────────────────┘  │  /claude    │
                              └─────────────┘"""
    _code_box(slide, arch, 0.5, 1.75, 12.3, 5.4)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 3  —  AI VISION
# ═══════════════════════════════════════════════════════════════════════════════

def s03_header(prs):
    _section_header(prs, 3, "AI & Vision", "ปัญญาประดิษฐ์และการวิเคราะห์ภาพ", ORANGE)


def s03_ai(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "What is AI?", "ปัญญาประดิษฐ์คืออะไร?")

    stages = [
        ("Rule-Based",   "กฎตายตัว\nIF วัตถุกลม AND สีแดง → THEN แอปเปิ้ล", GRAY),
        ("Machine\nLearning", "เรียนรู้จากตัวอย่าง\n1,000 รูป → เรียนรู้เอง → ทำนาย", CYAN),
        ("Generative AI", "สร้างคำตอบใหม่\nรับรูป + prompt → ตอบภาษาธรรมชาติ", ORANGE),
    ]
    for i, (title, desc, color) in enumerate(stages):
        left = 0.5 + i * 4.3
        _box(slide, left, 1.75, 4.0, 5.3, BG2, border=True)
        _box(slide, left, 1.75, 4.0, 0.55, color)
        _tb(slide, title, left + 0.15, 1.78, 3.7, 0.5,
            size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, desc, left + 0.15, 2.55, 3.7, 4.3,
            size=15, color=LIGHT, wrap=True)
        _tb(slide, "→" if i < 2 else "", left + 4.05, 3.8, 0.5, 0.5,
            size=24, color=CYAN, bold=True)


def s03_genai(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Generative AI", "AI ที่สร้างคำตอบใหม่จากข้อมูลที่ได้รับ")

    items = [
        ("📝  Text Generation",  "เขียนเรียงความ, สรุปข้อมูล, ตอบคำถาม"),
        ("🖼️  Image Generation", "สร้างภาพจากคำอธิบาย (Midjourney, DALL-E)"),
        ("👁️  Image Analysis",   "อธิบาย, จำแนก, นับ สิ่งที่อยู่ในภาพ  ← Workshop นี้!"),
        ("💻  Code Generation",  "เขียนโปรแกรมจาก requirements"),
    ]
    top = 1.85
    for icon_title, desc in items:
        highlight = "Workshop" in desc
        _box(slide, 0.5, top, 12.3, 1.05,
             RGBColor(0x1A, 0x2F, 0x1A) if highlight else BG2,
             border=highlight)
        _tb(slide, icon_title, 0.75, top + 0.1, 4.5, 0.5,
            size=19, color=CYAN if not highlight else GREEN, bold=True)
        _tb(slide, desc, 5.4, top + 0.1, 7.2, 0.7,
            size=16, color=LIGHT)
        top += 1.22


def s03_vision(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Vision AI", "AI ที่ 'มองเห็น' และเข้าใจภาพ")

    _tb(slide,
        "Vision AI รับภาพ + prompt แล้วตอบกลับเป็นข้อความ",
        0.5, 1.72, 12.5, 0.5, size=20, color=CYAN)

    capabilities = [
        ("🔍 Object Detection",     "ระบุวัตถุในภาพ: แมว, โต๊ะ, ขวด"),
        ("🔢 Counting",              "นับจำนวน: มีหนังสือ 3 เล่ม"),
        ("🎨 Color / Scene",        "วิเคราะห์สีและบรรยากาศ"),
        ("📊 Classification",        "จำแนกประเภท: ขยะรีไซเคิล"),
        ("💚 Health Assessment",    "ประเมินสุขภาพ: ต้นไม้มีใบเหลือง"),
        ("📝 Text Reading (OCR)",   "อ่านตัวหนังสือในภาพ"),
    ]
    top = 2.38
    for i, (cap, example) in enumerate(capabilities):
        col = 0 if i % 2 == 0 else 6.6
        row = top + (i // 2) * 1.38
        _box(slide, 0.5 + col, row, 6.3, 1.15, BG2, border=True)
        _tb(slide, cap, 0.75 + col, row + 0.08, 5.8, 0.4,
            size=17, color=ORANGE, bold=True)
        _tb(slide, example, 0.75 + col, row + 0.52, 5.8, 0.45,
            size=14, color=LIGHT)


def s03_gemini(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Gemini Vision (Google)", "Vision AI โดย Google — แนะนำสำหรับ Workshop")

    _box(slide, 0.5, 1.75, 12.3, 5.4, BG2, border=True)
    info = [
        ("Model",        "gemini-2.5-flash"),
        ("Provider",     "Google AI Studio"),
        ("API Key",      "ฟรีสำหรับนักเรียน — aistudio.google.com"),
        ("Free Tier",    "1,500 requests/วัน  (เพียงพอสำหรับ Workshop)"),
        ("SDK",          "google-genai  (python)"),
        ("ENV",          "VISION_PROVIDER=gemini"),
    ]
    top = 1.95
    for label, val in info:
        _tb(slide, label, 0.8, top, 2.5, 0.48,
            size=17, color=CYAN, bold=True)
        _tb(slide, val, 3.4, top, 9.2, 0.48,
            size=17, color=LIGHT)
        top += 0.75

    _tb(slide, "✅  แนะนำให้เริ่มด้วย Gemini เพราะ free tier สูงที่สุด",
        0.8, 6.6, 12.0, 0.45, size=14, color=GREEN, bold=True)


def s03_claude(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Claude Vision (Anthropic)", "Vision AI โดย Anthropic")

    _box(slide, 0.5, 1.75, 12.3, 5.4, BG2, border=True)
    info = [
        ("Model",     "claude-haiku-4-5"),
        ("Provider",  "Anthropic"),
        ("API Key",   "console.anthropic.com  (ต้องใส่บัตรเครดิต)"),
        ("Pricing",   "$1.00 / M input tokens  •  $5.00 / M output tokens"),
        ("SDK",       "anthropic  (python)"),
        ("ENV",       "VISION_PROVIDER=claude"),
    ]
    top = 1.95
    for label, val in info:
        _tb(slide, label, 0.8, top, 2.5, 0.48,
            size=17, color=ORANGE, bold=True)
        _tb(slide, val, 3.4, top, 9.2, 0.48,
            size=17, color=LIGHT)
        top += 0.75

    _tb(slide, "💡  ใช้ได้เช่นกัน — เปลี่ยนแค่ VISION_PROVIDER=claude ใน .env",
        0.8, 6.6, 12.0, 0.45, size=14, color=CYAN)


def s03_prompt(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Prompt Engineering", "วิธีเขียนคำสั่งให้ AI ตอบได้ตรงใจ")

    # Formula
    _box(slide, 0.5, 1.72, 12.3, 0.72, RGBColor(0x00, 0x28, 0x40), border=True)
    _tb(slide, "[Role]  +  [Task]  +  [Format]  +  [Language]",
        0.7, 1.77, 12.0, 0.6,
        size=22, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    parts = [
        ("Role",     "คุณคือผู้เชี่ยวชาญด้านสิ่งแวดล้อม",   "บอก AI ว่าเป็นใคร"),
        ("Task",     "วิเคราะห์ขยะในภาพและจำแนกประเภท",     "สิ่งที่ต้องทำ"),
        ("Format",   "ตอบเป็น bullet point",                   "รูปแบบคำตอบ"),
        ("Language", "ภาษาไทย",                               "ภาษาที่ใช้"),
    ]
    top = 2.65
    for label, example, desc in parts:
        _box(slide, 0.5, top, 12.3, 1.02, BG2, border=True)
        _tb(slide, label, 0.75, top + 0.06, 1.8, 0.4,
            size=17, color=ORANGE, bold=True)
        _tb(slide, f'"{example}"', 2.65, top + 0.06, 5.5, 0.4,
            size=15, color=GREEN)
        _tb(slide, desc, 8.3, top + 0.06, 4.3, 0.4,
            size=13, color=GRAY, italic=True)
        top += 1.18


def s03_workflow(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "AI Vision Workflow", "ข้อมูลไหลอย่างไร จากกล้องถึง AI")

    workflow = """\
  👤 นักเรียน เปิด Web UI (localhost:8000)
        │
        │  กรอก cam_url + prompt → กด "Capture & Analyze"
        ▼
  ⚙️  FastAPI Backend  (main.py)
        │
        ├─── GET http://<esp32-ip>/capture  ──► 📷 ESP32-CAM
        │         (JPEG image, ~10–50 KB)         └── ถ่ายภาพ → ส่งกลับ
        │
        ├─── image_bytes + prompt
        │         ▼
        │    Vision Provider  (VISION_PROVIDER env var)
        │         ├── mock     → canned response (ไม่ต้อง key)
        │         ├── gemini   → Google AI API
        │         └── claude   → Anthropic API
        │
        └── JSON Response  ──► 🌐 Web UI แสดงผล"""

    _code_box(slide, workflow, 0.5, 1.75, 12.3, 5.4)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 4  —  LAB 01
# ═══════════════════════════════════════════════════════════════════════════════

def s04_header(prs):
    slide = _new_slide(prs)
    _section_ghost(slide, 1)
    _box(slide, 0, 0, 0.12, 7.5, ORANGE)
    _box(slide, 0.12, 6.8, 13.21, 0.06, ORANGE)
    _tb(slide, "Lab 01", 4.5, 1.6, 8.5, 1.1, size=58, color=ORANGE, bold=True)
    _tb(slide, "ESP32-CAM Camera Server", 4.5, 2.7, 8.5, 0.85,
        size=30, color=WHITE, bold=True)
    _tb(slide, "ติดตั้งกล้อง • Flash Firmware • ดู Live Stream", 4.5, 3.55, 8.5, 0.6,
        size=18, color=LIGHT)
    _tb(slide, "เวลา: ~60 นาที", 4.5, 4.3, 4.0, 0.4, size=15, color=GRAY)


def s04_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Lab 01 — Overview", "สิ่งที่เราจะทำใน Lab นี้")

    steps = [
        "ต่อสาย ESP32-CAM กับ ESP32-CAM-MB board",
        "ติดตั้ง Arduino IDE และ ESP32 Board Package",
        "แก้ไข secrets.h ใส่ Wi-Fi SSID และ Password",
        "Flash firmware เข้า ESP32-CAM",
        "ทดสอบ Live Stream และ /capture endpoint",
    ]
    _checklist(slide, steps, top=1.85, size=19)

    _tb(slide, "🎯  เป้าหมาย: กล้อง online และ /capture ส่ง JPEG ได้",
        0.5, 6.7, 12.3, 0.45, size=15, color=GREEN, bold=True)


def s04_step1(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 1 — Wire ESP32-CAM")
    _step_badge(slide, 1, "ต่อสาย ESP32-CAM + MB Board")

    diagram = """\
  ESP32-CAM
  ┌──────────┐         ESP32-CAM-MB
  │          │◄──────► (เสียบให้แน่น ทิศทางถูก)
  │  GPIO0   │── ต่อลง GND ← ⚠️ ทำก่อน Upload
  │  5V      │── 5V (จาก MB Board)
  │  GND     │── GND
  └──────────┘
       │  USB Micro-B
       ▼
  ต่อกับ Notebook / PC"""
    _code_box(slide, diagram, 0.5, 2.5, 7.5, 4.5)

    _multiline(slide, [
        ("⚠️  GPIO0 → GND", ORANGE, 18, True),
        ("ต้องทำก่อน Upload ทุกครั้ง", LIGHT, 16, False),
        ("หลัง Upload เสร็จ:", LIGHT, 16, False),
        ("→ ถอด jumper GPIO0 ออก", CYAN, 16, False),
        ("→ กดปุ่ม RST บน MB Board", CYAN, 16, False),
    ], 8.2, 2.5, 4.9, 4.5, gap=8)


def s04_step2(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 2 — Arduino IDE Setup")
    _step_badge(slide, 2, "ตั้งค่า Arduino IDE")

    steps_txt = """\
1. เปิด Arduino IDE 2.x
2. File → Preferences
   → Additional boards manager URLs:
   https://dl.espressif.com/dl/package_esp32_index.json

3. Tools → Board → Boards Manager
   → ค้นหา "esp32" → Install (Espressif Systems)

4. Tools → Board → ESP32 Arduino
   → เลือก "AI Thinker ESP32-CAM"

5. Tools → Port → เลือก COMx (ของ MB Board)"""
    _code_box(slide, steps_txt, 0.5, 2.42, 12.3, 4.6)


def s04_step3(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 3 — Edit secrets.h")
    _step_badge(slide, 3, "ใส่ Wi-Fi credentials")

    _tb(slide, "คัดลอก  secrets.example.h  →  secrets.h  แล้วแก้ไข:",
        0.5, 2.42, 12.5, 0.45, size=16, color=LIGHT)

    code = """\
// secrets.h
#define WIFI_SSID  "ชื่อ Wi-Fi ของคุณ"
#define WIFI_PASS  "รหัส Wi-Fi ของคุณ"
"""
    _code_box(slide, code, 0.5, 3.0, 12.3, 1.6)

    _multiline(slide, [
        ("⚠️  secrets.h อยู่ใน .gitignore — ไม่ถูก commit ขึ้น GitHub", ORANGE, 15, True),
        ("✅  ใช้ 2.4 GHz เท่านั้น (ESP32 ไม่รองรับ 5 GHz)", GREEN, 15, False),
        ("✅  ตรวจสอบให้ครั้งสุดท้ายก่อน Upload", GREEN, 15, False),
    ], 0.5, 4.75, 12.5, 2.2, gap=10)


def s04_step4(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 4 — Flash Firmware")
    _step_badge(slide, 4, "Upload Firmware เข้า ESP32-CAM")

    steps_txt = """\
1. ตรวจสอบ: GPIO0 ต่อลง GND แล้ว
2. กด Upload ใน Arduino IDE (→)
3. รอ "Connecting......____"
   → ถ้าค้าง: ลองกด RST บน MB Board สั้นๆ
4. เห็น "Writing at 0x..." = กำลัง flash ✅
5. เห็น "Hard resetting..." = flash เสร็จ ✅
6. ถอด jumper GPIO0 ออก
7. กด RST อีกครั้ง → ESP32-CAM บูตปกติ"""
    _code_box(slide, steps_txt, 0.5, 2.42, 7.5, 4.6)

    _multiline(slide, [
        ("ปัญหาที่พบบ่อย:", ORANGE, 17, True),
        ("❌ ค้างที่ Connecting", LIGHT, 15, False),
        ("→ GPIO0 ไม่ได้ต่อ GND", GRAY, 14, False),
        ("❌ Upload failed", LIGHT, 15, False),
        ("→ เลือก Port ผิด", GRAY, 14, False),
        ("❌ ไม่เห็น Port", LIGHT, 15, False),
        ("→ ไม่มีไดรเวอร์ CH340", GRAY, 14, False),
    ], 8.3, 2.42, 4.8, 4.6, gap=6)


def s04_step5(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 1)
    _slide_title(slide, "Step 5 — Find IP & Test")
    _step_badge(slide, 5, "หา IP Address และทดสอบ")

    code = """\
# เปิด Serial Monitor (Tools → Serial Monitor)
# Baud Rate: 115200

# ควรเห็น:
Connecting to WiFi......
WiFi connected!
IP Address: 192.168.x.x   ← จด IP นี้ไว้!
Camera Ready! Stream: http://192.168.x.x:81/stream"""
    _code_box(slide, code, 0.5, 2.42, 7.5, 3.3)

    _multiline(slide, [
        ("ทดสอบใน Browser:", CYAN, 17, True),
        ("http://<IP>", LIGHT, 15, False),
        ("→ หน้า Web UI", GRAY, 14, False),
        ("http://<IP>/capture", LIGHT, 15, False),
        ("→ ดาวน์โหลด JPEG", GRAY, 14, False),
        ("http://<IP>:81/stream", LIGHT, 15, False),
        ("→ Live video stream", GRAY, 14, False),
    ], 8.3, 2.42, 4.8, 3.3, gap=8)

    _tb(slide, "📝  จด IP Address ไว้ใช้ใน Lab 02-03-04",
        0.5, 5.9, 12.3, 0.45, size=15, color=ORANGE, bold=True)


def s04_result(prs):
    slide = _new_slide(prs)
    _top_bar(slide, GREEN)
    _accent_bar(slide, color=GREEN)
    _lab_badge(slide, 1, GREEN)
    _slide_title(slide, "Lab 01 — Expected Result")

    _checklist(slide, [
        "ESP32-CAM เชื่อมต่อ Wi-Fi สำเร็จ",
        "Serial Monitor แสดง IP Address",
        "http://<IP>  เปิดได้ใน Browser",
        "http://<IP>/capture  ส่ง JPEG กลับมา",
        "http://<IP>:81/stream  แสดง Live Video",
    ], top=1.85)

    _tb(slide, "🎉  Lab 01 สำเร็จ! พร้อมสำหรับ Lab 02",
        0.5, 6.7, 12.3, 0.45, size=17, color=GREEN, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 5  —  LAB 02
# ═══════════════════════════════════════════════════════════════════════════════

def s05_header(prs):
    slide = _new_slide(prs)
    _section_ghost(slide, 2)
    _box(slide, 0, 0, 0.12, 7.5, CYAN)
    _box(slide, 0.12, 6.8, 13.21, 0.06, CYAN)
    _tb(slide, "Lab 02", 4.5, 1.6, 8.5, 1.1, size=58, color=CYAN, bold=True)
    _tb(slide, "REST API & Backend", 4.5, 2.7, 8.5, 0.85,
        size=30, color=WHITE, bold=True)
    _tb(slide, "รัน Backend • ทดสอบ API • Mock Provider", 4.5, 3.55, 8.5, 0.6,
        size=18, color=LIGHT)
    _tb(slide, "เวลา: ~60 นาที", 4.5, 4.3, 4.0, 0.4, size=15, color=GRAY)


def s05_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Lab 02 — Overview", "สิ่งที่เราจะทำใน Lab นี้")

    steps = [
        "สร้าง Python Virtual Environment (.venv)",
        "ติดตั้ง dependencies ด้วย pip",
        "ตั้งค่า .env ไฟล์  (VISION_PROVIDER=mock)",
        "รัน FastAPI server ด้วย uvicorn",
        "ทดสอบ upload รูปและดูผลผ่าน Web UI",
    ]
    _checklist(slide, steps, top=1.85, size=19)
    _tb(slide, "🎯  เป้าหมาย: เข้าใจ API data flow — ไม่ต้องใช้ API key",
        0.5, 6.7, 12.3, 0.45, size=15, color=GREEN, bold=True)


def s05_step1(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Step 1 — Create Virtual Environment")
    _step_badge(slide, 1, "สร้าง Python Virtual Environment", CYAN)

    code_win = """\
# Windows (PowerShell / Command Prompt)
cd code\\server
python -m venv .venv
.venv\\Scripts\\activate

# เห็น (.venv) หน้า prompt = สำเร็จ ✅"""

    code_mac = """\
# macOS / Linux
cd code/server
python3 -m venv .venv
source .venv/bin/activate

# เห็น (.venv) หน้า prompt = สำเร็จ ✅"""

    _tb(slide, "Windows:", 0.5, 2.42, 6.3, 0.35, size=15, color=CYAN, bold=True)
    _code_box(slide, code_win, 0.5, 2.8, 6.3, 2.8)
    _tb(slide, "macOS / Linux:", 7.0, 2.42, 6.1, 0.35, size=15, color=CYAN, bold=True)
    _code_box(slide, code_mac, 7.0, 2.8, 6.1, 2.8)

    _tb(slide, "💡  Virtual Environment แยก Python packages ของ project ออกจากระบบ",
        0.5, 5.75, 12.5, 0.45, size=14, color=GRAY, italic=True)


def s05_step2(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Step 2 — Install Dependencies")
    _step_badge(slide, 2, "ติดตั้ง Python Packages", CYAN)

    code = """\
pip install -r requirements.txt

# packages ที่จะถูกติดตั้ง:
#  fastapi          → web framework
#  uvicorn          → ASGI server
#  httpx            → fetch รูปจาก ESP32-CAM
#  python-dotenv    → อ่านค่าจาก .env
#  python-multipart → รับไฟล์ upload
#  pillow           → image processing
#  google-genai     → Gemini Vision SDK
#  anthropic        → Claude Vision SDK"""
    _code_box(slide, code, 0.5, 2.42, 12.3, 4.6)


def s05_step3(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Step 3 — Configure .env")
    _step_badge(slide, 3, "ตั้งค่าไฟล์ Environment", CYAN)

    _tb(slide, "คัดลอก  .env.example  →  .env  แล้วแก้ไข:",
        0.5, 2.42, 12.5, 0.4, size=16, color=LIGHT)

    code = """\
# .env  (Lab 02 — ใช้ mock provider ก่อน)
VISION_PROVIDER=mock

# เติมทีหลังใน Lab 03:
GEMINI_API_KEY=your_key_here
GEMINI_MODEL=gemini-2.5-flash
CAM_URL=http://192.168.x.x"""
    _code_box(slide, code, 0.5, 2.98, 12.3, 2.8)

    _multiline(slide, [
        ("⚠️  .env อยู่ใน .gitignore — ไม่ถูก commit ขึ้น GitHub", ORANGE, 15, True),
        ("✅  mock provider ไม่ต้องใช้ API key — ตอบ canned response", GREEN, 15, False),
    ], 0.5, 5.95, 12.5, 1.0, gap=8)


def s05_step4(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _lab_badge(slide, 2, CYAN)
    _slide_title(slide, "Step 4 — Run FastAPI Server")
    _step_badge(slide, 4, "รัน Backend Server", CYAN)

    code = """\
# ใน terminal (ที่ activate .venv แล้ว)
uvicorn main:app --reload --host 0.0.0.0

# ควรเห็น:
INFO:     Uvicorn running on http://0.0.0.0:8000
INFO:     Application startup complete.

# เปิดใน Browser:
http://localhost:8000          → Web UI
http://localhost:8000/health   → { "status": "ok", "vision_provider": "mock" }
http://localhost:8000/prompts  → รายการ preset prompts"""
    _code_box(slide, code, 0.5, 2.42, 12.3, 4.6)


def s05_result(prs):
    slide = _new_slide(prs)
    _top_bar(slide, GREEN)
    _accent_bar(slide, color=GREEN)
    _lab_badge(slide, 2, GREEN)
    _slide_title(slide, "Lab 02 — Expected Result")

    _checklist(slide, [
        "pip install สำเร็จ ไม่มี error",
        ".env ไฟล์สร้างแล้ว  VISION_PROVIDER=mock",
        "uvicorn รันได้ที่  http://localhost:8000",
        "Web UI เปิดได้และมี preset prompts dropdown",
        "Upload รูปทดสอบ → ได้รับ mock response กลับมา",
        "/health แสดง  vision_provider: mock",
    ], top=1.85)

    _tb(slide, "🎉  Lab 02 สำเร็จ! พร้อม connect AI จริงใน Lab 03",
        0.5, 6.7, 12.3, 0.45, size=17, color=GREEN, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 6  —  LAB 03
# ═══════════════════════════════════════════════════════════════════════════════

def s06_header(prs):
    slide = _new_slide(prs)
    _section_ghost(slide, 3)
    _box(slide, 0, 0, 0.12, 7.5, ORANGE)
    _box(slide, 0.12, 6.8, 13.21, 0.06, ORANGE)
    _tb(slide, "Lab 03", 4.5, 1.6, 8.5, 1.1, size=58, color=ORANGE, bold=True)
    _tb(slide, "AI Vision Analysis", 4.5, 2.7, 8.5, 0.85,
        size=30, color=WHITE, bold=True)
    _tb(slide, "API Key จริง • ภาพจากกล้อง • Prompt Engineering", 4.5, 3.55, 8.5, 0.6,
        size=18, color=LIGHT)
    _tb(slide, "เวลา: ~75 นาที", 4.5, 4.3, 4.0, 0.4, size=15, color=GRAY)


def s06_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Lab 03 — Overview", "เชื่อม Vision AI จริง กับกล้อง ESP32-CAM")

    steps = [
        "สมัครและรับ Gemini API Key จาก Google AI Studio",
        "ตั้งค่า GEMINI_API_KEY และ VISION_PROVIDER=gemini ใน .env",
        "Restart uvicorn server",
        "ชี้กล้องและกด Capture & Analyze ครั้งแรก",
        "ทดลอง prompt อย่างน้อย 3 แบบ สังเกตความแตกต่าง",
    ]
    _checklist(slide, steps, top=1.85, size=19)
    _tb(slide, "🎯  เป้าหมาย: AI จริงวิเคราะห์ภาพจากกล้อง ESP32-CAM",
        0.5, 6.7, 12.3, 0.45, size=15, color=GREEN, bold=True)


def s06_step1(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Step 1 — Get Gemini API Key")
    _step_badge(slide, 1, "รับ API Key จาก Google AI Studio")

    steps_txt = """\
1. เปิด Browser → aistudio.google.com
2. Sign in ด้วย Google Account
3. คลิก  "Get API key"
4. คลิก  "Create API key"
5. คัดลอก key (ขึ้นต้นด้วย AIza...)
6. วางใน .env:
   GEMINI_API_KEY=AIzaXXXXXXXXXXXXXXXXXX"""
    _code_box(slide, steps_txt, 0.5, 2.42, 7.5, 4.0)

    _multiline(slide, [
        ("Free Tier:", CYAN, 17, True),
        ("1,500 req / วัน", GREEN, 15, False),
        ("ไม่ต้องใช้บัตรเครดิต", GREEN, 15, False),
        ("", GRAY, 14, False),
        ("⚠️  อย่าแชร์ API key", ORANGE, 15, True),
        ("อย่า commit ลง GitHub", ORANGE, 15, False),
    ], 8.3, 2.42, 4.8, 4.0, gap=8)


def s06_step2(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Step 2 — Update .env & Restart")
    _step_badge(slide, 2, "แก้ .env แล้ว Restart Server")

    code = """\
# .env  (Lab 03)
VISION_PROVIDER=gemini
GEMINI_API_KEY=AIzaXXXXXXXXXXXXXXXXXX
GEMINI_MODEL=gemini-2.5-flash
CAM_URL=http://192.168.x.x   ← IP จาก Lab 01

# Restart uvicorn:
# กด Ctrl+C แล้วรันใหม่
uvicorn main:app --reload --host 0.0.0.0

# ตรวจสอบ:
# GET http://localhost:8000/health
# → { "vision_provider": "gemini" }"""
    _code_box(slide, code, 0.5, 2.42, 12.3, 4.6)


def s06_step3(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Step 3 — First AI Analysis!")
    _step_badge(slide, 3, "วิเคราะห์ภาพครั้งแรกด้วย Gemini")

    code = """\
1. เปิด http://localhost:8000
2. กรอก ESP32-CAM URL:  http://192.168.x.x
3. เลือก Prompt: "อธิบายสิ่งที่เห็นในภาพนี้"
4. กดปุ่ม  [ Capture & Analyze ]
5. รอ 2-5 วินาที...
6. ดูผลลัพธ์ใน Result panel

# ถ้าได้ error:
# 401 → API key ผิด / ไม่ได้ตั้ง
# 429 → quota เต็ม → รอ 1 นาที หรือ mock
# 502 → network / camera problem"""
    _code_box(slide, code, 0.5, 2.42, 12.3, 4.6)


def s06_step4(prs):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _lab_badge(slide, 3)
    _slide_title(slide, "Step 4 — Prompt Experiments")
    _step_badge(slide, 4, "ทดลอง Prompt อย่างน้อย 3 แบบ")

    prompts = [
        ("Prompt A (กว้าง)",    "อธิบายสิ่งที่เห็นในภาพนี้"),
        ("Prompt B (เจาะจง)",   "นับจำนวนวัตถุแต่ละประเภทในภาพ ตอบเป็น list"),
        ("Prompt C (มี Role)",  "คุณคือผู้เชี่ยวชาญ วิเคราะห์ภาพและให้คะแนน 1-10"),
    ]
    top = 2.1
    for label, prompt in prompts:
        _box(slide, 0.5, top, 12.3, 1.3, BG2, border=True)
        _tb(slide, label, 0.75, top + 0.08, 3.5, 0.4,
            size=16, color=ORANGE, bold=True)
        _tb(slide, f'"{prompt}"', 0.75, top + 0.52, 11.8, 0.55,
            size=15, color=GREEN)
        top += 1.5

    _tb(slide, "สังเกต: Prompt ที่ต่างกัน → คำตอบที่ต่างกัน",
        0.5, 6.7, 12.3, 0.45, size=15, color=CYAN, bold=True)


def s06_result(prs):
    slide = _new_slide(prs)
    _top_bar(slide, GREEN)
    _accent_bar(slide, color=GREEN)
    _lab_badge(slide, 3, GREEN)
    _slide_title(slide, "Lab 03 — Expected Result")

    _checklist(slide, [
        "Gemini API Key ตั้งค่าใน .env สำเร็จ",
        "/health แสดง  vision_provider: gemini",
        "กด Capture & Analyze แล้วได้ผลจาก Gemini จริง",
        "ทดลองอย่างน้อย 3 prompt แตกต่างกัน",
        "เห็นว่า prompt ส่งผลต่อความละเอียดของคำตอบ",
    ], top=1.85)

    _tb(slide, "🎉  Lab 03 สำเร็จ! พร้อมสร้าง Mini Project ใน Lab 04",
        0.5, 6.7, 12.3, 0.45, size=17, color=GREEN, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 7  —  LAB 04
# ═══════════════════════════════════════════════════════════════════════════════

def s07_header(prs):
    slide = _new_slide(prs)
    _section_ghost(slide, 4)
    _box(slide, 0, 0, 0.12, 7.5, CYAN)
    _box(slide, 0.12, 6.8, 13.21, 0.06, CYAN)
    _tb(slide, "Lab 04", 4.5, 1.6, 8.5, 1.1, size=58, color=CYAN, bold=True)
    _tb(slide, "AI Vision Mini Project", 4.5, 2.7, 8.5, 0.85,
        size=30, color=WHITE, bold=True)
    _tb(slide, "เลือก Project • เขียน Prompt • สร้าง AI ของตัวเอง", 4.5, 3.55, 8.5, 0.6,
        size=18, color=LIGHT)
    _tb(slide, "เวลา: ~75 นาที", 4.5, 4.3, 4.0, 0.4, size=15, color=GRAY)


def s07_overview(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _lab_badge(slide, 4, CYAN)
    _slide_title(slide, "Lab 04 — Overview", "สร้าง AI Vision Project ของตัวเอง")

    flow = [
        ("1. เลือกหัวข้อ",    "เลือก 1 ใน 5 preset หรือคิดขึ้นเอง"),
        ("2. เขียน Prompt",   "ใช้สูตร Role + Task + Format + Language"),
        ("3. ทดสอบ",          "กด Capture & Analyze ≥ 3 ครั้ง"),
        ("4. ปรับปรุง",        "แก้ prompt จนได้ผลที่พอใจ"),
        ("5. นำเสนอ",         "demo + อธิบาย prompt + สิ่งที่เรียนรู้ (2 นาที)"),
    ]
    top = 1.75
    for step, desc in flow:
        _box(slide, 0.5, top, 12.3, 0.95, BG2, border=True)
        _tb(slide, step, 0.75, top + 0.1, 2.5, 0.45,
            size=18, color=CYAN, bold=True)
        _tb(slide, desc, 3.4, top + 0.1, 9.2, 0.6,
            size=16, color=LIGHT)
        top += 1.1


def s07_step1(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _lab_badge(slide, 4, CYAN)
    _slide_title(slide, "Step 1 — Choose Your Project")
    _step_badge(slide, 1, "เลือกหัวข้อ Mini Project", CYAN)

    projects = [
        ("🗑️  Waste Classification",       "จำแนกขยะเป็น รีไซเคิล / อินทรีย์ / อันตราย"),
        ("🌿  Plant Health Analysis",       "ประเมินสุขภาพต้นไม้จากสี ใบ และลักษณะ"),
        ("🔢  Object Counting",             "นับจำนวนวัตถุแต่ละประเภทในภาพ"),
        ("🖥️  Desk Cleanliness",           "ให้คะแนนความสะอาด / เป็นระเบียบของโต๊ะ"),
        ("🔧  Equipment Detection",         "ตรวจสอบอุปกรณ์ในห้องปฏิบัติการ"),
    ]
    top = 2.42
    for icon_name, desc in projects:
        _tb(slide, icon_name, 0.75, top, 4.5, 0.42, size=17, color=ORANGE, bold=True)
        _tb(slide, desc, 5.4, top, 7.7, 0.42, size=15, color=LIGHT)
        top += 0.75

    _tb(slide, "💡  หรือคิดหัวข้อขึ้นเองก็ได้!",
        0.5, 6.75, 12.3, 0.42, size=15, color=CYAN, italic=True)


def s07_step2(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _lab_badge(slide, 4, CYAN)
    _slide_title(slide, "Step 2 — Write Your Prompt")
    _step_badge(slide, 2, "เขียน Custom Prompt สำหรับ Project", CYAN)

    code = """\
# สูตร: Role + Task + Format + Language

คุณคือ [บทบาทของ AI]

วิเคราะห์ภาพนี้และ[สิ่งที่ต้องทำ]

ตอบในรูปแบบ:
[หัวข้อ 1]: [ค่า]
[หัวข้อ 2]: [ค่า]
สรุป: [ประโยคสรุป]

ตอบเป็นภาษาไทย"""
    _code_box(slide, code, 0.5, 2.45, 7.5, 4.5)

    _multiline(slide, [
        ("Tips:", CYAN, 17, True),
        ("• ทดสอบอย่างน้อย 3 ครั้ง", LIGHT, 15, False),
        ("• ถ้า AI ตอบกว้างเกิน", LIGHT, 15, False),
        ("  → เพิ่มรายละเอียดใน Format", GRAY, 14, False),
        ("• ถ้า AI ตอบไม่ตรง", LIGHT, 15, False),
        ("  → ปรับ Role ให้ชัดขึ้น", GRAY, 14, False),
        ("• ดูตัวอย่างใน", LIGHT, 15, False),
        ("  docs/Prompt-Library.md", GREEN, 14, False),
    ], 8.2, 2.45, 4.9, 4.5, gap=8)


def s07_tips(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _lab_badge(slide, 4, CYAN)
    _slide_title(slide, "Presentation Tips", "วิธีนำเสนอ Mini Project (2 นาที)")

    items = [
        ("🎯  Demo Live",    "ชี้กล้องไปที่วัตถุจริง แล้วกด Capture & Analyze"),
        ("📝  แสดง Prompt",  "บอกว่าเขียน prompt อะไร และทำไมถึงเลือกแบบนั้น"),
        ("📊  แสดงผลลัพธ์", "AI ตอบอะไร? ตรงกับที่คาดหวังไหม?"),
        ("💡  สิ่งที่เรียนรู้", "ได้เรียนรู้อะไรจาก workshop นี้?"),
    ]
    top = 1.9
    for icon_title, desc in items:
        _box(slide, 0.5, top, 12.3, 1.15, BG2, border=True)
        _tb(slide, icon_title, 0.75, top + 0.1, 3.0, 0.5,
            size=19, color=CYAN, bold=True)
        _tb(slide, desc, 3.9, top + 0.1, 9.0, 0.7, size=16, color=LIGHT)
        top += 1.32


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 8  —  MINI PROJECTS
# ═══════════════════════════════════════════════════════════════════════════════

def s08_header(prs):
    _section_header(prs, 8, "Mini Projects", "ไอเดียโปรเจกต์สำหรับ Lab 04", ORANGE)


def _project_slide(prs, icon, title_en, title_th, prompt, output):
    slide = _new_slide(prs)
    _top_bar(slide, ORANGE)
    _accent_bar(slide, color=ORANGE)
    _tb(slide, f"{icon}  {title_en}", 0.5, 0.2, 12.5, 0.7,
        size=26, color=ORANGE, bold=True)
    _tb(slide, title_th, 0.5, 0.85, 12.5, 0.45,
        size=16, color=LIGHT)
    _box(slide, 0.5, 1.42, 12.3, 0.05, ORANGE)

    _tb(slide, "Prompt:", 0.5, 1.6, 2.0, 0.4, size=15, color=CYAN, bold=True)
    _code_box(slide, prompt, 0.5, 2.0, 12.3, 2.5)

    _tb(slide, "Expected Output:", 0.5, 4.65, 3.5, 0.4, size=15, color=GREEN, bold=True)
    _code_box(slide, output, 0.5, 5.05, 12.3, 2.1)


def s08_waste(prs):
    _project_slide(prs,
        "🗑️", "Waste Classification", "จำแนกประเภทขยะ",
        'คุณคือระบบ AI จำแนกขยะอัจฉริยะ\nวิเคราะห์ขยะในภาพและตอบในรูปแบบนี้:\n🗑️ ประเภทขยะ: [ชื่อ]\n♻️ หมวดหมู่: รีไซเคิล / อินทรีย์ / ทั่วไป / อันตราย\n✅ วิธีทิ้ง: [คำแนะนำ]\nตอบภาษาไทย',
        '🗑️ ประเภทขยะ: ขวดพลาสติก\n♻️ หมวดหมู่: รีไซเคิล\n✅ วิธีทิ้ง: ทิ้งในถังขยะรีไซเคิลสีเหลือง ล้างให้สะอาดก่อนทิ้ง')


def s08_plant(prs):
    _project_slide(prs,
        "🌿", "Plant Health Analysis", "วิเคราะห์สุขภาพต้นไม้",
        'คุณคือนักพฤกษศาสตร์ผู้เชี่ยวชาญ\nวิเคราะห์ต้นไม้ในภาพและตอบ:\n🌿 ชนิด: [ถ้าระบุได้]\n💚 สุขภาพ: ดีมาก / ดี / ปานกลาง / ต้องการดูแลด่วน\n🔍 สังเกต: [สีใบ, จุด, ใบเหี่ยว]\n💧 การรดน้ำ: [แนะนำ]\n🌱 คำแนะนำ: [2-3 ข้อ]',
        '🌿 ชนิด: ต้นมอนสเตอร่า\n💚 สุขภาพ: ดี\n🔍 สังเกต: ใบสีเขียวเข้ม สมบูรณ์\n💧 การรดน้ำ: เพียงพอแล้ว\n🌱 คำแนะนำ: วางในที่มีแสงรำไร')


def s08_count(prs):
    _project_slide(prs,
        "🔢", "Object Counting", "นับจำนวนวัตถุ",
        'นับและจำแนกวัตถุทุกชิ้นในภาพนี้อย่างละเอียด\nตอบในรูปแบบ:\n[ชื่อวัตถุ]: [จำนวน]\n...\nรวมทั้งหมด: [จำนวน] ชิ้น\nหมายเหตุ: [วัตถุที่มองไม่ชัด]',
        'โน้ตบุ๊ค: 1\nกาแฟ: 1\nหนังสือ: 3\nปากกา: 2\nแอปเปิ้ล: 1\nรวมทั้งหมด: 8 ชิ้น')


def s08_desk(prs):
    _project_slide(prs,
        "🖥️", "Desk Cleanliness Inspection", "ประเมินความเป็นระเบียบของโต๊ะ",
        'คุณคือผู้ตรวจสอบความสะอาด ประเมินโต๊ะนี้:\n⭐ คะแนน: [1-5 ดาว]\n📦 สิ่งของ: [รายการ]\n✅ จุดที่ดี: [สิ่งที่เป็นระเบียบ]\n⚠️ ปรับปรุง: [สิ่งที่ควรจัดใหม่]\n💡 คำแนะนำ: [2-3 ข้อ]',
        '⭐ คะแนน: 4 ดาว\n📦 สิ่งของ: โน้ตบุ๊ค, เมาส์, แก้วน้ำ, สมุด\n✅ จุดที่ดี: จัดสายไฟเป็นระเบียบดี\n⚠️ ปรับปรุง: มีเศษกระดาษบนโต๊ะ')


def s08_equipment(prs):
    _project_slide(prs,
        "🔧", "Equipment Detection", "ตรวจสอบอุปกรณ์",
        'คุณคือระบบตรวจสอบอุปกรณ์ห้องปฏิบัติการ\nวิเคราะห์ภาพและระบุ:\n🔧 อุปกรณ์ที่พบ: [รายการ]\n✅ พร้อมใช้งาน: [รายการ]\n⚠️ ต้องตรวจสอบ: [รายการ]\n❌ หายหรือชำรุด: [ถ้ามี]\nสรุป: [1 ประโยค]',
        '🔧 พบ: ESP32-CAM, สาย USB, breadboard\n✅ พร้อม: ESP32-CAM, สาย USB\n⚠️ ตรวจสอบ: breadboard (มีสีออกซิไดซ์)\n❌ หาย: -\nสรุป: อุปกรณ์ส่วนใหญ่พร้อมใช้งาน')


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 9  —  CLOSING
# ═══════════════════════════════════════════════════════════════════════════════

def s09_careers(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Future Career Paths", "เส้นทางอาชีพที่เกี่ยวข้อง")

    careers = [
        ("🤖  AI Engineer",          "พัฒนาโมเดล AI, fine-tuning, deployment"),
        ("📡  IoT Developer",         "ออกแบบระบบเซนเซอร์, firmware, cloud"),
        ("📊  Data Scientist",        "วิเคราะห์ข้อมูล, ML models, visualization"),
        ("🤖  Robotics Engineer",    "หุ่นยนต์ที่ใช้ computer vision นำทาง"),
        ("🛡️  AI Safety Researcher", "ความปลอดภัยและจริยธรรมของ AI"),
        ("🚀  AI Product Manager",   "นำทีมสร้าง product ที่ขับเคลื่อนด้วย AI"),
    ]
    top = 1.72
    for i, (title, desc) in enumerate(careers):
        col = 0 if i % 2 == 0 else 6.6
        row = top + (i // 2) * 1.6
        _box(slide, 0.5 + col, row, 6.3, 1.38, BG2, border=True)
        _tb(slide, title, 0.75 + col, row + 0.1, 5.8, 0.45,
            size=17, color=CYAN, bold=True)
        _tb(slide, desc, 0.75 + col, row + 0.58, 5.8, 0.55,
            size=13, color=LIGHT)


def s09_resources(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Additional Resources", "แหล่งเรียนรู้เพิ่มเติม")

    resources = [
        ("📖  Workshop GitHub",       "github.com/Pachara966/esp32cam-ai-vision-workshop"),
        ("🤖  Google AI Studio",      "aistudio.google.com  (Gemini API Key)"),
        ("🧠  Anthropic Console",     "console.anthropic.com  (Claude API Key)"),
        ("📚  FastAPI Docs",          "fastapi.tiangolo.com"),
        ("📷  ESP32 Arduino Docs",    "docs.espressif.com"),
        ("🎓  Google ML Crash Course","developers.google.com/machine-learning/crash-course"),
    ]
    top = 1.72
    for icon_title, url in resources:
        _box(slide, 0.5, top, 12.3, 0.88, BG2, border=True)
        _tb(slide, icon_title, 0.75, top + 0.08, 3.5, 0.4,
            size=16, color=ORANGE, bold=True)
        _tb(slide, url, 4.4, top + 0.08, 8.2, 0.55,
            size=14, color=CYAN)
        top += 1.02


def s09_summary(prs):
    slide = _new_slide(prs)
    _top_bar(slide)
    _accent_bar(slide)
    _slide_title(slide, "Workshop Summary", "สรุป Workshop วันนี้")

    summary = [
        ("Lab 01", "ติดตั้ง ESP32-CAM + Flash Firmware + Live Stream",      ORANGE),
        ("Lab 02", "รัน FastAPI Backend + ทำความเข้าใจ REST API",             CYAN),
        ("Lab 03", "เชื่อม Gemini Vision + Prompt Engineering",              ORANGE),
        ("Lab 04", "สร้าง AI Vision Mini Project ของตัวเอง",                 GREEN),
    ]
    top = 1.8
    for lab, desc, color in summary:
        _box(slide, 0.5, top, 12.3, 1.05, BG2, border=True)
        _box(slide, 0.5, top, 1.5, 1.05, color)
        _tb(slide, lab, 0.5, top + 0.2, 1.5, 0.55,
            size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, desc, 2.2, top + 0.22, 10.5, 0.55,
            size=17, color=LIGHT)
        top += 1.22

    _tb(slide, "🎉  ภายใน 1 วัน นักเรียนสร้างระบบ AI Vision ที่ใช้งานได้จริง!",
        0.5, 6.65, 12.3, 0.5, size=17, color=GREEN, bold=True, align=PP_ALIGN.CENTER)


def s09_qa(prs):
    slide = _new_slide(prs)
    _box(slide, 0, 0, 13.33, 7.5, BG)   # plain dark bg
    _box(slide, 0, 0, 13.33, 0.09, CYAN)
    _box(slide, 0, 7.41, 13.33, 0.09, CYAN)

    _tb(slide, "?", 2.5, 0.5, 8.33, 5.8,
        size=260, color=GHOST, bold=True, align=PP_ALIGN.CENTER)

    _tb(slide, "Q&A", 0.5, 2.0, 12.33, 2.5,
        size=100, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _tb(slide, "ถามได้เลยครับ / Feel free to ask!",
        0.5, 5.8, 12.33, 0.6,
        size=22, color=CYAN, align=PP_ALIGN.CENTER)


def s09_thankyou(prs):
    slide = _new_slide(prs)
    _box(slide, 0, 0, 13.33, 7.5, BG)
    _box(slide, 0, 0, 13.33, 0.12, ORANGE)
    _box(slide, 0, 7.38, 13.33, 0.12, ORANGE)
    _box(slide, 0, 0.12, 0.1, 7.26, CYAN)

    _tb(slide, "ขอบคุณทุกคน!", 0.5, 0.9, 12.33, 1.5,
        size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _tb(slide, "Thank You!", 0.5, 2.3, 12.33, 1.0,
        size=40, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    _box(slide, 3.5, 3.45, 6.33, 0.06, GRAY)

    _tb(slide, "มหาวิทยาลัยเทคโนโลยีราชมงคลธัญบุรี (RMUTT)",
        0.5, 3.65, 12.33, 0.55,
        size=19, color=LIGHT, align=PP_ALIGN.CENTER)

    _tb(slide, "[ชื่อผู้สอน / Instructor Name]",
        0.5, 4.3, 12.33, 0.45,
        size=16, color=GRAY, align=PP_ALIGN.CENTER)

    _tb(slide, "github.com/Pachara966/esp32cam-ai-vision-workshop",
        0.5, 5.0, 12.33, 0.45,
        size=15, color=CYAN, align=PP_ALIGN.CENTER)

    _tb(slide, "#ESP32CAM_AI_RMUTT",
        0.5, 5.55, 12.33, 0.45,
        size=18, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

    _tb(slide, "🎉  สร้างโดย: ESP32-CAM AI Vision Workshop Team",
        0.5, 6.55, 12.33, 0.45,
        size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Building slides...")

    # Section 0 — Opening
    s00_cover(prs);     print("  ✅ 01 Cover")
    s00_overview(prs);  print("  ✅ 02 Overview")
    s00_outcomes(prs);  print("  ✅ 03 Learning Outcomes")
    s00_agenda(prs);    print("  ✅ 04 Agenda")

    # Section 1 — IoT & Hardware
    s01_header(prs);   print("  ✅ 05 Section 1 Header")
    s01_iot(prs);      print("  ✅ 06 What is IoT")
    s01_esp32cam(prs); print("  ✅ 07 ESP32-CAM Intro")
    s01_hardware(prs); print("  ✅ 08 Hardware Tour")
    s01_wifi(prs);     print("  ✅ 09 Wi-Fi Basics")

    # Section 2 — REST API
    s02_header(prs);   print("  ✅ 10 Section 2 Header")
    s02_api(prs);      print("  ✅ 11 What is REST API")
    s02_http(prs);     print("  ✅ 12 HTTP GET vs POST")
    s02_json(prs);     print("  ✅ 13 JSON Basics")
    s02_fastapi(prs);  print("  ✅ 14 FastAPI Overview")

    # Section 3 — AI Vision
    s03_header(prs);    print("  ✅ 15 Section 3 Header")
    s03_ai(prs);        print("  ✅ 16 What is AI")
    s03_genai(prs);     print("  ✅ 17 Generative AI")
    s03_vision(prs);    print("  ✅ 18 Vision AI")
    s03_gemini(prs);    print("  ✅ 19 Gemini Vision")
    s03_claude(prs);    print("  ✅ 20 Claude Vision")
    s03_prompt(prs);    print("  ✅ 21 Prompt Engineering")
    s03_workflow(prs);  print("  ✅ 22 AI Vision Workflow")

    # Section 4 — Lab 01
    s04_header(prs);   print("  ✅ 23 Lab 01 Header")
    s04_overview(prs); print("  ✅ 24 Lab 01 Overview")
    s04_step1(prs);    print("  ✅ 25 Lab 01 Step 1")
    s04_step2(prs);    print("  ✅ 26 Lab 01 Step 2")
    s04_step3(prs);    print("  ✅ 27 Lab 01 Step 3")
    s04_step4(prs);    print("  ✅ 28 Lab 01 Step 4")
    s04_step5(prs);    print("  ✅ 29 Lab 01 Step 5")
    s04_result(prs);   print("  ✅ 30 Lab 01 Result")

    # Section 5 — Lab 02
    s05_header(prs);   print("  ✅ 31 Lab 02 Header")
    s05_overview(prs); print("  ✅ 32 Lab 02 Overview")
    s05_step1(prs);    print("  ✅ 33 Lab 02 Step 1")
    s05_step2(prs);    print("  ✅ 34 Lab 02 Step 2")
    s05_step3(prs);    print("  ✅ 35 Lab 02 Step 3")
    s05_step4(prs);    print("  ✅ 36 Lab 02 Step 4")
    s05_result(prs);   print("  ✅ 37 Lab 02 Result")

    # Section 6 — Lab 03
    s06_header(prs);   print("  ✅ 38 Lab 03 Header")
    s06_overview(prs); print("  ✅ 39 Lab 03 Overview")
    s06_step1(prs);    print("  ✅ 40 Lab 03 Step 1")
    s06_step2(prs);    print("  ✅ 41 Lab 03 Step 2")
    s06_step3(prs);    print("  ✅ 42 Lab 03 Step 3")
    s06_step4(prs);    print("  ✅ 43 Lab 03 Step 4")
    s06_result(prs);   print("  ✅ 44 Lab 03 Result")

    # Section 7 — Lab 04
    s07_header(prs);   print("  ✅ 45 Lab 04 Header")
    s07_overview(prs); print("  ✅ 46 Lab 04 Overview")
    s07_step1(prs);    print("  ✅ 47 Lab 04 Step 1")
    s07_step2(prs);    print("  ✅ 48 Lab 04 Step 2")
    s07_tips(prs);     print("  ✅ 49 Lab 04 Tips")

    # Section 8 — Mini Projects
    s08_header(prs);     print("  ✅ 50 Mini Projects Header")
    s08_waste(prs);      print("  ✅ 51 Waste Classification")
    s08_plant(prs);      print("  ✅ 52 Plant Health")
    s08_count(prs);      print("  ✅ 53 Object Counting")
    s08_desk(prs);       print("  ✅ 54 Desk Cleanliness")
    s08_equipment(prs);  print("  ✅ 55 Equipment Detection")

    # Section 9 — Closing
    s09_careers(prs);   print("  ✅ 56 Future Careers")
    s09_resources(prs); print("  ✅ 57 Resources")
    s09_summary(prs);   print("  ✅ 58 Summary")
    s09_qa(prs);        print("  ✅ 59 Q&A")
    s09_thankyou(prs);  print("  ✅ 60 Thank You")

    out = "Workshop.pptx"
    prs.save(out)
    print(f"\n🎉  Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
